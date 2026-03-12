"""
render_from_template.py
-----------------------
Stencil-based slide renderer for PPMaker.

Rather than drawing shapes from scratch, this module clones ("duplicates") a
reference slide from the stencil PPTX, then surgically replaces text content
while preserving all original formatting (font, size, color, bold, spacing).

Main public surface:
    StencilRenderer(stencil_pptx_path, catalog, brand_tokens)
    renderer.render(outline, visual_director=None, logo_path=None) -> BytesIO

Internal helpers (also exported for testing):
    duplicate_slide_in_prs(prs, slide_idx) -> Slide
    classify_text_zones(slide)             -> list[dict]
    set_para_text(para, text)
    replace_text_preserving_format(tf, new_text, max_chars)
    replace_bullets_preserving_format(tf, bullets, max_items)
"""

import copy
import io
import logging
import os
from typing import Optional

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

from models.schemas import StorytellingOutline

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Safe PPTX opener (handles OneDrive-locked files)
# ---------------------------------------------------------------------------

def _open_pptx_safe(path: str) -> Presentation:
    """
    Open a PPTX file even when it is locked by OneDrive/cloud sync.
    Strategy:
      1. Try direct open(path, 'rb').
      2. On PermissionError, look for a pre-existing copy in TEMP
         (named ambev.pptx or ambev_inspect.pptx — created by bash cp).
      3. As last resort, raise so the caller can fall back gracefully.
    """
    import tempfile

    try:
        with open(path, "rb") as fh:
            return Presentation(io.BytesIO(fh.read()))
    except PermissionError:
        pass

    # Look for known temp copies (created by the bash cp commands used in
    # tools/inspect_template.py or manual setup)
    tmp_dir = tempfile.gettempdir()
    for candidate_name in ("ambev.pptx", "ambev_inspect.pptx",
                           os.path.basename(path)):
        candidate = os.path.join(tmp_dir, candidate_name)
        if os.path.isfile(candidate):
            logger.info("Stencil locked by OneDrive — using temp copy: %s", candidate)
            try:
                with open(candidate, "rb") as fh:
                    return Presentation(io.BytesIO(fh.read()))
            except Exception:
                continue

    raise PermissionError(
        f"Cannot open stencil PPTX (locked by OneDrive/sync): {path}\n"
        "Fix: run  cp '<stencil.pptx>' $TEMP/ambev.pptx  once to create a temp copy."
    )


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SLIDE_W_IN = 13.33
SLIDE_H_IN = 7.50
SLIDE_AREA = SLIDE_W_IN * SLIDE_H_IN

# Relationship namespace (used when remapping rIds)
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Hardcoded fallback stencil map: layout hint → 0-based slide index
FALLBACK_STENCIL_MAP: dict[str, int] = {
    "title":              4,
    "hero":               7,
    "content":            1,
    "two-column":        29,
    "cards":             15,
    "dashboard":          9,
    "chart-placeholder":  9,
    "image-text":        10,
    "closing":           45,
}

# Typical "lorem ipsum" fragments used in template placeholder text
_LOREM_FRAGMENTS = (
    "lorem", "ipsum", "dolor", "consectetur", "adipiscing",
    "pellentesque", "viverra", "facilisis",
)

import re as _re


def _extract_key_metric(talking_points: list) -> str:
    """Extract a short numeric KPI/metric string from talking points.

    Examples:
        'Volume cresceu +12% YoY' → '+12%'
        'NPS atingiu 78 pontos'   → '78'
        'Eficiência operacional'  → ''   (no number found)
    """
    for pt in talking_points:
        m = _re.search(r"[+\-]?\d+(?:[.,]\d+)?\s*(?:%|pts?|k\b|x\b)?", pt, _re.I)
        if m:
            return m.group(0).strip()
    return ""


# ---------------------------------------------------------------------------
# Slide duplication
# ---------------------------------------------------------------------------

def duplicate_slide_in_prs(prs: Presentation, slide_idx: int):
    """
    Clone slide[slide_idx] and append the clone to prs.

    Shares image/media parts with the source slide (same PPTX file).
    Remaps relationship IDs so charts, images, and hyperlinks resolve
    correctly on the new slide.

    Returns the new pptx.slide.Slide object.
    """
    src    = prs.slides[slide_idx]
    layout = src.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # ---- Replace spTree content (all shapes) ----
    sp_tree_new = new_slide.shapes._spTree
    sp_tree_src = src.shapes._spTree
    sp_tree_new.clear()
    for child in sp_tree_src:
        sp_tree_new.append(copy.deepcopy(child))

    # ---- Replicate relationships ----
    rId_map: dict[str, str] = {}
    for rId, rel in src.part.rels.items():
        try:
            if rel.is_external:
                new_rId = new_slide.part.add_relationship(
                    rel.reltype, rel.target_ref, is_external=True
                ).rId
            else:
                new_rId = new_slide.part.add_relationship(
                    rel.reltype, rel.target_part
                ).rId
            rId_map[rId] = new_rId
        except Exception as exc:
            logger.debug("Could not replicate rel %s: %s", rId, exc)

    # ---- Remap rId attribute values in cloned XML ----
    for elem in sp_tree_new.iter():
        for attr, val in list(elem.attrib.items()):
            if val in rId_map:
                elem.set(attr, rId_map[val])

    return new_slide


# ---------------------------------------------------------------------------
# Text-zone classification
# ---------------------------------------------------------------------------

def classify_text_zones(slide) -> list[dict]:
    """
    Walk all shapes on a slide (including groups) and return a list of
    text-zone descriptors sorted by (top, left) position.

    Each descriptor:
        {
            "shape": <Shape>,
            "role":  "TITLE|SUBTITLE|HIGHLIGHT|BODY|LABEL|CAPTION",
            "top":   float,   # inches
            "area":  float,   # sq inches
        }
    """
    zones: list[dict] = []

    def _process(shape):
        try:
            stype = shape.shape_type
        except Exception:
            stype = None

        if stype == MSO_SHAPE_TYPE.GROUP:
            try:
                for child in shape.shapes:
                    _process(child)
            except Exception:
                pass
            return

        if not hasattr(shape, "text_frame"):
            return

        try:
            text = shape.text_frame.text.strip()
            if not text:
                return
        except Exception:
            return

        try:
            left   = shape.left   / 914400  # EMU → inches
            top    = shape.top    / 914400
            width  = shape.width  / 914400
            height = shape.height / 914400
        except Exception:
            left = top = width = height = 0.0

        area = width * height

        # Determine max font size in shape
        # Check run-level first, then paragraph-level (theme-inherited fonts
        # are often only visible at the paragraph level, not on individual runs)
        max_pt = 0.0
        try:
            for para in shape.text_frame.paragraphs:
                # Run-level size
                for run in para.runs:
                    try:
                        sz = run.font.size
                        if sz:
                            pt = sz / 12700
                            if pt > max_pt:
                                max_pt = pt
                    except Exception:
                        pass
                # Paragraph-level size (fallback when run size is None)
                if max_pt == 0.0:
                    try:
                        sz = para.font.size
                        if sz:
                            pt = sz / 12700
                            if pt > max_pt:
                                max_pt = pt
                    except Exception:
                        pass
        except Exception:
            pass

        # Count non-empty paragraphs
        try:
            n_paras = len([p for p in shape.text_frame.paragraphs if p.text.strip()])
        except Exception:
            n_paras = 1

        zones.append({
            "shape":    shape,
            "role":     None,    # filled in second pass below
            "top":      top,
            "left":     left,
            "area":     area,
            "max_pt":   max_pt,
            "n_paras":  n_paras,
            "width":    width,
            "height":   height,
        })

    for shape in slide.shapes:
        _process(shape)

    if not zones:
        return zones

    zones.sort(key=lambda z: (z["top"], z.get("left", 0)))

    # --- Second pass: assign roles ---
    # Step 1: classify shapes that have explicit font sizes using font-size rules.
    # Step 2: for shapes with max_pt == 0 (theme-inherited, unreadable from XML),
    #         use positional heuristics so TITLE + BODY are still identified.

    title_assigned = False
    body_assigned  = False
    subtitle_assigned = False

    for z in zones:
        max_pt  = z["max_pt"]
        top     = z["top"]
        area    = z["area"]
        n_paras = z["n_paras"]
        width   = z["width"]
        if max_pt == 0.0:
            continue  # handled in second sweep below
        if max_pt >= 60:
            z["role"] = "HIGHLIGHT"
        elif max_pt >= 24 and top < 2.0:
            z["role"] = "TITLE"
            title_assigned = True
        elif max_pt >= 14 and top < 3.0:
            z["role"] = "SUBTITLE"
            subtitle_assigned = True
        elif (area / SLIDE_AREA) > 0.20 and n_paras > 1:
            z["role"] = "BODY"
            body_assigned = True
        elif width < 3.0:
            z["role"] = "LABEL"
        else:
            z["role"] = "CAPTION"

    # Collect unclassified 0pt zones and apply positional heuristics
    zero_zones = [z for z in zones if z["role"] is None]
    if zero_zones:
        # Among 0pt shapes, pick TITLE = topmost with area > 5% (if none yet)
        # and BODY = largest area with n_paras > 1 (if none yet)
        # Body candidate: multi-para shape, OR large single-para shape (>15% area)
        body_cands = sorted(
            [z for z in zero_zones
             if z["n_paras"] > 1 or (z["area"] / SLIDE_AREA) > 0.15],
            key=lambda z: -z["area"],
        )
        zero_body = body_cands[0] if body_cands and not body_assigned else None

        # TITLE candidate: topmost shape with area > 4% that isn't the body
        title_cands = sorted(
            [z for z in zero_zones if z["area"] / SLIDE_AREA > 0.04 and z is not zero_body],
            key=lambda z: z["top"],
        )
        zero_title = title_cands[0] if title_cands and not title_assigned else None

        # If still no body candidate (all shapes have area < 15%, n_paras = 1),
        # pick the largest remaining 0pt shape that isn't TITLE as BODY.
        if zero_body is None and not body_assigned:
            remaining = sorted(
                [z for z in zero_zones if z is not zero_title and
                 z["area"] / SLIDE_AREA > 0.03],
                key=lambda z: -z["area"],
            )
            zero_body = remaining[0] if remaining else None

        for z in zero_zones:
            if z is zero_title:
                z["role"] = "TITLE"
            elif z is zero_body:
                z["role"] = "BODY"
            else:
                text = z["shape"].text_frame.text.strip()
                if len(text) <= 15 and z["top"] < 3.5:
                    z["role"] = "HIGHLIGHT"
                elif z["width"] < 3.0:
                    z["role"] = "LABEL"
                else:
                    z["role"] = "CAPTION"

    # Ensure every zone has a role (safety net)
    for z in zones:
        if z["role"] is None:
            z["role"] = "CAPTION"

    return zones


# ---------------------------------------------------------------------------
# Text replacement helpers
# ---------------------------------------------------------------------------

def set_para_text(para, text: str) -> None:
    """
    Set the text of a paragraph to `text`, preserving the first run's
    formatting. Removes all runs beyond the first.
    """
    try:
        runs = para.runs
        if not runs:
            run = para.add_run()
            run.text = text
            return

        # Keep first run, clear the rest
        runs[0].text = text
        for run in runs[1:]:
            r_elem = run._r
            if r_elem is not None:
                r_elem.getparent().remove(r_elem)
    except Exception as exc:
        logger.debug("set_para_text error: %s", exc)


def replace_text_preserving_format(
    tf,
    new_text: str,
    max_chars: int = 500,
) -> None:
    """
    Replace ALL text in a text_frame with new_text.

    - Preserves formatting (font, size, color, bold) of the first existing run.
    - Handles overflow: truncates with "…" if len(new_text) > max_chars.
    - Handles multi-line: if new_text contains "\\n", creates multiple paragraphs.
    - Reduces font size in 2pt steps if text is very long (floor 8pt).
    """
    if not new_text:
        return

    # Truncate if needed
    if len(new_text) > max_chars:
        new_text = new_text[:max_chars - 1] + "…"

    # Extract formatting template from the first paragraph / run
    template_para = None
    template_run  = None
    try:
        for para in tf.paragraphs:
            if para.runs:
                template_para = para
                template_run  = para.runs[0]
                break
    except Exception:
        pass

    # Capture first-run formatting (copy values, not references)
    tmpl_font_name  = None
    tmpl_font_size  = None
    tmpl_bold       = None
    tmpl_color_rgb  = None
    tmpl_alignment  = None

    if template_run:
        try:
            tmpl_font_name = template_run.font.name
        except Exception:
            pass
        try:
            tmpl_font_size = template_run.font.size  # in EMU (12700 per pt)
        except Exception:
            pass
        try:
            tmpl_bold = template_run.font.bold
        except Exception:
            pass
        try:
            if template_run.font.color and template_run.font.color.type:
                tmpl_color_rgb = template_run.font.color.rgb
        except Exception:
            pass

    if template_para:
        try:
            tmpl_alignment = template_para.alignment
        except Exception:
            pass

    # ---- Font size adaptation based on text length ----
    adjusted_size = tmpl_font_size
    if tmpl_font_size:
        char_len = len(new_text)
        if char_len > 200:
            adjusted_size = max(int(8 * 12700), tmpl_font_size - int(4 * 12700))
        elif char_len > 120:
            adjusted_size = max(int(8 * 12700), tmpl_font_size - int(2 * 12700))

    # ---- Clear all existing paragraphs except the first ----
    try:
        txBody = tf._txBody
        paras  = txBody.findall(qn("a:p"))
        for p_elem in paras[1:]:
            txBody.remove(p_elem)
    except Exception as exc:
        logger.debug("Could not clear extra paragraphs: %s", exc)

    # ---- Write lines ----
    lines = new_text.split("\n")

    for line_idx, line in enumerate(lines):
        if line_idx == 0:
            # Use the existing first paragraph
            try:
                first_para = tf.paragraphs[0]
                set_para_text(first_para, line)
                # Re-apply formatting to the first (now only) run
                _apply_run_format(first_para, tmpl_font_name, adjusted_size, tmpl_bold, tmpl_color_rgb)
                if tmpl_alignment is not None:
                    first_para.alignment = tmpl_alignment
            except Exception as exc:
                logger.debug("Error writing first line: %s", exc)
        else:
            # Add a new paragraph cloned from the template
            try:
                new_para = tf.add_paragraph()
                new_para.text = line
                _apply_run_format(new_para, tmpl_font_name, adjusted_size, tmpl_bold, tmpl_color_rgb)
                if tmpl_alignment is not None:
                    new_para.alignment = tmpl_alignment
            except Exception as exc:
                logger.debug("Error writing line %d: %s", line_idx, exc)


def replace_bullets_preserving_format(
    tf,
    bullets: list[str],
    max_items: int = 8,
) -> None:
    """
    Replace text_frame content with a bullet list.

    - Preserves the formatting of the first existing paragraph as template.
    - Adds "• " prefix if original template text started with "•" or "-".
    - Limits to max_items, appending "…" if truncated.
    """
    if not bullets:
        return

    # Detect whether template uses a bullet prefix
    use_prefix = False
    try:
        first_text = tf.paragraphs[0].text.strip()
        if first_text.startswith("•") or first_text.startswith("-"):
            use_prefix = True
    except Exception:
        pass

    # Capture template formatting
    tmpl_font_name = None
    tmpl_font_size = None
    tmpl_bold      = None
    tmpl_color_rgb = None
    tmpl_alignment = None

    try:
        first_para = tf.paragraphs[0]
        if first_para.runs:
            run = first_para.runs[0]
            try:
                tmpl_font_name = run.font.name
            except Exception:
                pass
            try:
                tmpl_font_size = run.font.size
            except Exception:
                pass
            try:
                tmpl_bold = run.font.bold
            except Exception:
                pass
            try:
                if run.font.color and run.font.color.type:
                    tmpl_color_rgb = run.font.color.rgb
            except Exception:
                pass
        try:
            tmpl_alignment = first_para.alignment
        except Exception:
            pass
    except Exception:
        pass

    # Truncate bullet list
    display_bullets = list(bullets[:max_items])
    if len(bullets) > max_items:
        display_bullets.append("…")

    # ---- Clear existing paragraphs except first ----
    try:
        txBody = tf._txBody
        paras  = txBody.findall(qn("a:p"))
        for p_elem in paras[1:]:
            txBody.remove(p_elem)
    except Exception as exc:
        logger.debug("Could not clear bullet paragraphs: %s", exc)

    for idx, bullet_text in enumerate(display_bullets):
        display_text = ("• " + bullet_text) if use_prefix else bullet_text

        if idx == 0:
            try:
                first_para = tf.paragraphs[0]
                set_para_text(first_para, display_text)
                _apply_run_format(first_para, tmpl_font_name, tmpl_font_size, tmpl_bold, tmpl_color_rgb)
                if tmpl_alignment is not None:
                    first_para.alignment = tmpl_alignment
            except Exception as exc:
                logger.debug("Error setting bullet 0: %s", exc)
        else:
            try:
                new_para = tf.add_paragraph()
                new_para.text = display_text
                _apply_run_format(new_para, tmpl_font_name, tmpl_font_size, tmpl_bold, tmpl_color_rgb)
                if tmpl_alignment is not None:
                    new_para.alignment = tmpl_alignment
            except Exception as exc:
                logger.debug("Error setting bullet %d: %s", idx, exc)


def _apply_run_format(
    para,
    font_name: Optional[str],
    font_size,       # EMU int or None
    bold: Optional[bool],
    color_rgb,       # RGBColor or None
) -> None:
    """Apply formatting attributes to the first run of a paragraph."""
    try:
        runs = para.runs
        if not runs:
            return
        run = runs[0]
        if font_name:
            try:
                run.font.name = font_name
            except Exception:
                pass
        if font_size:
            try:
                run.font.size = font_size
            except Exception:
                pass
        if bold is not None:
            try:
                run.font.bold = bold
            except Exception:
                pass
        if color_rgb is not None:
            try:
                run.font.color.rgb = color_rgb
            except Exception:
                pass
    except Exception as exc:
        logger.debug("_apply_run_format error: %s", exc)


# ---------------------------------------------------------------------------
# StencilRenderer
# ---------------------------------------------------------------------------

class StencilRenderer:
    """
    Renders a StorytellingOutline into a PPTX by cloning stencil slides
    and replacing text content.

    Parameters
    ----------
    stencil_pptx_path : str
        Absolute or relative path to the stencil .pptx file.
        The file is opened with open(path, 'rb') to handle Windows paths
        with spaces.
    catalog : dict
        Parsed layout_catalog.json content.
    brand_tokens : dict
        Parsed brand_tokens.json content.
    """

    def __init__(
        self,
        stencil_pptx_path: str,
        catalog: dict,
        brand_tokens: dict,
    ) -> None:
        self._stencil_path  = stencil_pptx_path
        self._catalog_dict  = catalog
        self._brand_tokens  = brand_tokens

        # Build layout_id → catalog entry index for fast lookup
        self._entry_by_id: dict[str, dict] = {
            e["layout_id"]: e
            for e in catalog.get("catalog", [])
        }

    # ------------------------------------------------------------------
    # Public: render
    # ------------------------------------------------------------------

    def render(
        self,
        outline: StorytellingOutline,
        visual_director=None,
        logo_path: Optional[str] = None,
    ) -> io.BytesIO:
        """
        Render the entire deck.

        For each slide in the outline:
          1. Determine which stencil slide to clone.
          2. Clone it (duplicate_slide_in_prs).
          3. Replace text zones with slide content.
          4. Run QA checks.

        After all cloning, removes the original stencil slides so only
        the generated content remains.

        Returns
        -------
        io.BytesIO
            In-memory PPTX buffer (seeked to 0).
        """
        # Open the stencil — handle OneDrive-locked files via shell copy fallback
        prs = _open_pptx_safe(self._stencil_path)

        n_original = len(prs.slides)
        logger.info("Stencil loaded: %d original slides.", n_original)

        for slide_data in outline.slides:
            layout_hint = slide_data.layout

            # Resolve stencil index
            stencil_idx = None
            layout_entry: Optional[dict] = None

            if visual_director is not None:
                slide_content = {
                    "layout_hint":       layout_hint,
                    "title":             slide_data.title,
                    "bullets":           slide_data.talking_points,
                    "has_placeholder":   slide_data.has_placeholder,
                    "placeholder_hint":  slide_data.placeholder_hint,
                }
                deck_context = {
                    "deck_title":   outline.title,
                    "audience":     outline.audience,
                    "tone":         "professional",
                    "slide_index":  slide_data.index,
                    "total_slides": outline.total_slides,
                }
                layout_id   = visual_director.select_layout(slide_content, deck_context)
                stencil_idx = visual_director.get_stencil_index(layout_id)
                layout_entry = visual_director.get_catalog_entry(layout_id)

            if stencil_idx is None:
                stencil_idx = FALLBACK_STENCIL_MAP.get(layout_hint, 1)

            # Guard: clamp to valid range
            stencil_idx = max(0, min(stencil_idx, n_original - 1))

            logger.debug(
                "Slide %d (%s) → stencil idx %d",
                slide_data.index, layout_hint, stencil_idx,
            )

            try:
                new_slide = duplicate_slide_in_prs(prs, stencil_idx)
            except Exception as exc:
                logger.warning(
                    "Failed to duplicate stencil %d: %s. Using blank layout.",
                    stencil_idx, exc,
                )
                new_slide = prs.slides.add_slide(prs.slide_layouts[0])

            # Replace content
            try:
                self._replace_slide_content(new_slide, slide_data, layout_entry)
            except Exception as exc:
                logger.warning("Content replacement failed for slide %d: %s", slide_data.index, exc)

            # QA check
            warnings = self._qa_check(new_slide, slide_data)
            for w in warnings:
                logger.warning("QA [slide %d]: %s", slide_data.index, w)

        # ---- Remove original stencil slides ----
        _delete_first_n_slides(prs, n_original)
        logger.info("Removed %d original stencil slides.", n_original)
        logger.info("Final slide count: %d", len(prs.slides))

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf

    # ------------------------------------------------------------------
    # Content replacement
    # ------------------------------------------------------------------

    def _replace_slide_content(self, slide, slide_data, layout_entry: Optional[dict]) -> None:
        """
        Map slide_data fields onto the appropriate text zones of the cloned slide.

        Mapping rules:
          TITLE    → slide_data.title
          SUBTITLE → first talking point (if any)
          BODY     → all talking points as bullet list
          HIGHLIGHT → skip (keep decorative)
          LABEL    → skip (keep decorative numbers / category tags)
          CAPTION  → takeaway (first talking point or title)
        """
        zones = classify_text_zones(slide)

        if not zones:
            logger.debug("No text zones found on cloned slide — skipping replacement.")
            return

        title_done     = False
        subtitle_done  = False
        body_done      = False
        highlight_done = False
        caption_idx    = 0
        label_idx      = 0

        takeaway = slide_data.talking_points[0] if slide_data.talking_points else slide_data.title

        for zone in zones:
            shape = zone["shape"]
            role  = zone["role"]

            # HIGHLIGHT: replace with extracted key metric (first numeric value found)
            if role == "HIGHLIGHT":
                if not highlight_done:
                    metric = _extract_key_metric(slide_data.talking_points)
                    if metric:
                        try:
                            replace_text_preserving_format(shape.text_frame, metric, max_chars=30)
                            highlight_done = True
                        except Exception:
                            pass
                continue  # always skip remaining HIGHLIGHT zones

            # LABEL: replace with short sequential talking point content
            if role == "LABEL":
                pts = slide_data.talking_points
                if label_idx < len(pts):
                    try:
                        replace_text_preserving_format(
                            shape.text_frame, pts[label_idx][:60], max_chars=60
                        )
                    except Exception:
                        pass
                    label_idx += 1
                continue

            try:
                tf = shape.text_frame
            except Exception:
                continue

            if role == "TITLE" and not title_done:
                replace_text_preserving_format(tf, slide_data.title, max_chars=200)
                title_done = True

            elif role == "SUBTITLE" and not subtitle_done:
                sub = slide_data.talking_points[0] if slide_data.talking_points else ""
                if sub:
                    replace_text_preserving_format(tf, sub, max_chars=200)
                subtitle_done = True

            elif role == "BODY" and not body_done:
                if slide_data.talking_points:
                    replace_bullets_preserving_format(tf, slide_data.talking_points)
                body_done = True

            elif role == "CAPTION":
                # Distribute talking points across CAPTION zones sequentially
                pts = slide_data.talking_points
                text = pts[caption_idx] if caption_idx < len(pts) else takeaway
                caption_idx += 1
                replace_text_preserving_format(tf, text, max_chars=300)

        if not title_done:
            logger.warning("No TITLE zone found on slide; title '%s' not placed.", slide_data.title)
        if slide_data.talking_points and not body_done:
            logger.warning("No BODY zone found on slide; bullets not placed.")

    # ------------------------------------------------------------------
    # QA checks
    # ------------------------------------------------------------------

    def _qa_check(self, slide, slide_data) -> list[str]:
        """
        Run lightweight QA checks on the rendered slide.

        Returns a list of warning strings (empty = all OK).
        """
        warnings: list[str] = []

        # Check for un-replaced lorem ipsum placeholder text.
        # Only flag LARGE shapes (area > 3% of slide) to avoid spamming warnings
        # on decorative label/kpi shapes that legitimately keep template text.
        for shape in slide.shapes:
            try:
                if not hasattr(shape, "text_frame"):
                    continue
                text_lower = shape.text_frame.text.lower()
                if not any(frag in text_lower for frag in _LOREM_FRAGMENTS):
                    continue
                # Only warn for shapes large enough to be content zones
                area = (shape.width / 914400) * (shape.height / 914400)
                if area / SLIDE_AREA < 0.03:
                    continue  # skip small decorative shapes
                warnings.append(
                    f"ZONE_NOT_REPLACED: shape '{shape.name}' still contains placeholder text."
                )
            except Exception:
                pass

        # Check for out-of-bounds shapes
        for shape in slide.shapes:
            try:
                left = shape.left / 914400
                top  = shape.top  / 914400
                if left < -0.1 or top < -0.1:
                    warnings.append(
                        f"SHAPE_OUT_OF_BOUNDS: shape '{shape.name}' at ({left:.2f}, {top:.2f})."
                    )
            except Exception:
                pass

        # Title length check
        if len(slide_data.title) > 80:
            warnings.append(
                f"TITLE_TOO_LONG: title has {len(slide_data.title)} chars (max 80)."
            )

        # Too many bullets
        if len(slide_data.talking_points) > 8:
            warnings.append(
                f"TOO_MANY_BULLETS: {len(slide_data.talking_points)} bullets (max 8)."
            )

        return warnings


# ---------------------------------------------------------------------------
# Slide deletion helper
# ---------------------------------------------------------------------------

def _delete_first_n_slides(prs: Presentation, n: int) -> None:
    """
    Delete the first `n` slides from the presentation.

    Uses the rId from _sldIdLst to drop the relationship from the package part,
    then removes the element from the list.
    """
    for _ in range(n):
        if len(prs.slides._sldIdLst) == 0:
            break
        rId = prs.slides._sldIdLst[0].get(f"{{{_R_NS}}}id")
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception as exc:
                logger.debug("drop_rel %s: %s", rId, exc)
        del prs.slides._sldIdLst[0]
