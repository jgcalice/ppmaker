"""Tests for services/xml_renderer.py — XML-based PPTX renderer.

Covers the public contract:
    generate_pptx_xml(outline, template_path, template_meta) -> io.BytesIO

Tests follow AAA (Arrange / Act / Assert) and mock only at the boundary
(subprocess, os.path.exists) — never internal logic.
"""

import io
import subprocess
import sys
import tempfile
from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation as PptxPresentation


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def template_pptx_path():
    """Real global template — skip if not present."""
    path = Path(__file__).parent.parent.parent / "template_padrao" / "global" / "template-01.pptx"
    if not path.exists():
        pytest.skip("Template global template-01.pptx não disponível")
    return str(path)


@pytest.fixture
def small_outline_model(sample_outline_model):
    """Reuse conftest fixture — 3 slides, all layouts covered."""
    return sample_outline_model


# ---------------------------------------------------------------------------
# test_generate_pptx_xml_returns_bytesio
# ---------------------------------------------------------------------------

def test_generate_pptx_xml_returns_bytesio(
    small_outline_model, sample_template_meta, template_pptx_path
):
    """generate_pptx_xml must return an io.BytesIO instance at position 0."""
    # Arrange
    from services.xml_renderer import generate_pptx_xml

    # Act
    result = generate_pptx_xml(small_outline_model, template_pptx_path, sample_template_meta)

    # Assert
    assert isinstance(result, io.BytesIO), "Return type must be io.BytesIO"
    assert result.tell() == 0, "Buffer must be seeked to position 0"
    assert result.getbuffer().nbytes > 0, "Buffer must not be empty"


# ---------------------------------------------------------------------------
# test_generate_pptx_xml_produces_valid_pptx
# ---------------------------------------------------------------------------

def test_generate_pptx_xml_produces_valid_pptx(
    small_outline_model, sample_template_meta, template_pptx_path
):
    """The returned BytesIO must be openable as a valid PPTX with python-pptx."""
    # Arrange
    from services.xml_renderer import generate_pptx_xml

    # Act
    buffer = generate_pptx_xml(small_outline_model, template_pptx_path, sample_template_meta)

    # Assert — python-pptx raises if the ZIP/XML structure is corrupt
    prs = PptxPresentation(buffer)
    assert prs is not None, "python-pptx must be able to open the generated PPTX"


# ---------------------------------------------------------------------------
# test_generate_pptx_xml_slide_count_matches_outline
# ---------------------------------------------------------------------------

def test_generate_pptx_xml_slide_count_matches_outline(
    small_outline_model, sample_template_meta, template_pptx_path
):
    """Number of slides in the PPTX must equal outline.total_slides."""
    # Arrange
    from services.xml_renderer import generate_pptx_xml

    # Act
    buffer = generate_pptx_xml(small_outline_model, template_pptx_path, sample_template_meta)

    # Assert
    prs = PptxPresentation(buffer)
    assert len(prs.slides) == small_outline_model.total_slides, (
        f"Expected {small_outline_model.total_slides} slides, "
        f"got {len(prs.slides)}"
    )


# ---------------------------------------------------------------------------
# test_generate_pptx_xml_raises_when_scripts_missing
# ---------------------------------------------------------------------------

def test_generate_pptx_xml_raises_when_scripts_missing(
    small_outline_model, sample_template_meta, template_pptx_path
):
    """RuntimeError must be raised when required scripts are absent."""
    # Arrange — patch os.path.exists so all script paths return False
    from services import xml_renderer

    original_exists = Path.exists

    def _missing_scripts(self):
        # Only lie about files inside scripts/ — leave template alone
        if "scripts" in str(self) and self.suffix == ".py":
            return False
        return original_exists(self)

    with patch.object(Path, "exists", _missing_scripts):
        # Act & Assert
        with pytest.raises(RuntimeError, match="required scripts not found"):
            xml_renderer.generate_pptx_xml(
                small_outline_model, template_pptx_path, sample_template_meta
            )


# ---------------------------------------------------------------------------
# test_generate_pptx_xml_uses_temp_dir
# ---------------------------------------------------------------------------

def test_generate_pptx_xml_uses_temp_dir(
    small_outline_model, sample_template_meta, template_pptx_path
):
    """Renderer must use tempfile, not leave files in the working directory.

    After generate_pptx_xml returns, no 'output.pptx' or 'template.pptx'
    should exist in the current working directory.
    """
    # Arrange
    from services.xml_renderer import generate_pptx_xml
    import os

    cwd_before = set(os.listdir("."))

    # Act
    generate_pptx_xml(small_outline_model, template_pptx_path, sample_template_meta)

    # Assert — nothing new was created in cwd
    cwd_after = set(os.listdir("."))
    new_files = cwd_after - cwd_before
    pptx_artifacts = {f for f in new_files if f.endswith(".pptx") or f == "unpacked"}
    assert not pptx_artifacts, (
        f"Renderer left files in working directory: {pptx_artifacts}"
    )
