"""Tests for ported Anthropic scripts:
    - scripts/office/unpack.py  (unpack function)
    - scripts/office/pack.py    (pack function)
    - scripts/clean.py          (clean_unused_files function)
    - scripts/thumbnail.py      (generate_thumbnail_grid / convert_to_images)

Tests follow AAA, mock only at the boundary (subprocess, shutil.which),
and skip gracefully when optional dependencies are absent.
"""

import shutil
import sys
import tempfile
import zipfile
from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation as PptxPresentation


# ---------------------------------------------------------------------------
# sys.path helpers — scripts/ lives one level above backend/
# ---------------------------------------------------------------------------

_PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
_SCRIPTS_ROOT = _PROJECT_ROOT / "scripts"

if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

if str(_SCRIPTS_ROOT) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_ROOT))


# ---------------------------------------------------------------------------
# Shared fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def template_pptx_path():
    """Real global template — skip if not present."""
    path = _PROJECT_ROOT / "template_padrao" / "global" / "template-01.pptx"
    if not path.exists():
        pytest.skip("Template global template-01.pptx não disponível")
    return str(path)


@pytest.fixture
def no_soffice(monkeypatch):
    """Simulate absence of LibreOffice by making shutil.which return None for soffice/pdftoppm."""
    original_which = shutil.which

    def _fake_which(cmd):
        if cmd in ("soffice", "pdftoppm"):
            return None
        return original_which(cmd)

    monkeypatch.setattr("shutil.which", _fake_which)
    # Also patch inside thumbnail module (already imported)
    try:
        import scripts.thumbnail as _thumb
        monkeypatch.setattr(_thumb, "_soffice_available", lambda: False)
    except ImportError:
        pass


# ---------------------------------------------------------------------------
# test_unpack_creates_directory_structure
# ---------------------------------------------------------------------------

def test_unpack_creates_directory_structure(template_pptx_path, tmp_path):
    """unpack() must create a directory with PPTX structure: ppt/, ppt/slides/, [Content_Types].xml."""
    # Arrange
    from scripts.office.unpack import unpack

    output_dir = str(tmp_path / "unpacked")

    # Act
    _, message = unpack(template_pptx_path, output_dir)

    # Assert — no error and required structure is present
    assert not message.startswith("Error"), f"unpack returned error: {message}"

    unpacked_path = Path(output_dir)
    assert unpacked_path.is_dir(), "Output directory must be created"
    assert (unpacked_path / "[Content_Types].xml").exists(), "[Content_Types].xml must exist"
    assert (unpacked_path / "ppt").is_dir(), "ppt/ directory must exist"
    assert (unpacked_path / "ppt" / "slides").is_dir(), "ppt/slides/ directory must exist"


# ---------------------------------------------------------------------------
# test_pack_produces_valid_zip
# ---------------------------------------------------------------------------

def test_pack_produces_valid_zip(template_pptx_path, tmp_path):
    """pack() must produce a valid ZIP file with .pptx extension."""
    # Arrange
    from scripts.office.unpack import unpack
    from scripts.office.pack import pack

    unpacked_dir = str(tmp_path / "unpacked")
    output_pptx = str(tmp_path / "output.pptx")

    unpack(template_pptx_path, unpacked_dir)

    # Act
    _, message = pack(unpacked_dir, output_pptx, original_file=template_pptx_path, validate=True)

    # Assert
    assert not message.startswith("Error"), f"pack returned error: {message}"
    output_path = Path(output_pptx)
    assert output_path.exists(), "Output .pptx file must exist"
    assert zipfile.is_zipfile(str(output_path)), "Output file must be a valid ZIP"


# ---------------------------------------------------------------------------
# test_clean_removes_orphaned_files  (idempotency check)
# ---------------------------------------------------------------------------

def test_clean_removes_orphaned_files(template_pptx_path, tmp_path):
    """clean_unused_files must not crash on a clean directory (idempotent).

    On a freshly-unpacked template with all slides registered, running clean
    twice must not raise and the second run must remove nothing extra.
    """
    # Arrange
    from scripts.office.unpack import unpack
    from scripts.clean import clean_unused_files

    unpacked_dir = tmp_path / "unpacked"
    unpack(template_pptx_path, str(unpacked_dir))

    # Act — first clean (may remove genuinely orphaned items)
    removed_first = clean_unused_files(unpacked_dir)

    # Act — second clean (should be a no-op)
    removed_second = clean_unused_files(unpacked_dir)

    # Assert — second run is stable (no cascading removals)
    assert isinstance(removed_first, list), "clean_unused_files must return a list"
    assert removed_second == [], (
        f"Second clean should remove nothing, got: {removed_second}"
    )


# ---------------------------------------------------------------------------
# test_thumbnail_raises_without_soffice
# ---------------------------------------------------------------------------

def test_thumbnail_raises_without_soffice(template_pptx_path, tmp_path, no_soffice):
    """generate_thumbnail_grid must raise RuntimeError when soffice is absent."""
    # Arrange — no_soffice fixture patches shutil.which + _soffice_available
    from scripts.thumbnail import generate_thumbnail_grid

    output_prefix = str(tmp_path / "thumbnails")

    # Act & Assert
    with pytest.raises(RuntimeError, match="LibreOffice"):
        generate_thumbnail_grid(template_pptx_path, output_prefix)


# ---------------------------------------------------------------------------
# test_unpack_pack_roundtrip
# ---------------------------------------------------------------------------

def test_unpack_pack_roundtrip(template_pptx_path, tmp_path):
    """unpack + pack of a real template must yield an openable PPTX with the same slide count."""
    # Arrange
    from scripts.office.unpack import unpack
    from scripts.office.pack import pack

    unpacked_dir = str(tmp_path / "unpacked")
    output_pptx = str(tmp_path / "roundtrip.pptx")

    # Measure original slide count
    original_prs = PptxPresentation(template_pptx_path)
    original_slide_count = len(original_prs.slides)

    # Act
    _, unpack_msg = unpack(template_pptx_path, unpacked_dir)
    assert not unpack_msg.startswith("Error"), f"unpack failed: {unpack_msg}"

    _, pack_msg = pack(unpacked_dir, output_pptx, original_file=template_pptx_path, validate=True)
    assert not pack_msg.startswith("Error"), f"pack failed: {pack_msg}"

    # Assert — result is a valid PPTX with same slide count
    roundtrip_prs = PptxPresentation(output_pptx)
    assert len(roundtrip_prs.slides) == original_slide_count, (
        f"Slide count mismatch after roundtrip: "
        f"original={original_slide_count}, roundtrip={len(roundtrip_prs.slides)}"
    )
