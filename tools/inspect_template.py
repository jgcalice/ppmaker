"""
inspect_template.py
-------------------
Standalone reverse-engineering script for PPMaker.

Reads every .pptx in template_padrao/local/example/, inspects all shapes,
and produces two JSON files in the tools/ directory:

  - layout_catalog.json  : per-slide feature catalogue with text/image zones
  - brand_tokens.json    : font, color and typography tokens for the brand

Run from the ppmaker/ root:
    python tools/inspect_template.py

Dependencies: python-pptx, lxml  (both installed with python-pptx)
"""

import io
import json
import os
import sys
import glob
import traceback
from collections import Counter, defaultdict
from datetime import datetime, timezone
from lxml import etree

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SLIDE_W_IN = 13.33
SLIDE_H_IN = 7.50
SLIDE_W_EMU = int(SLIDE_W_IN * 914400)
SLIDE_H_EMU = int(SLIDE_H_IN * 914400)
SLIDE_AREA = SLIDE_W_IN * SLIDE_H_IN  # ~99.975 sq in

# Known AB InBev brand colours (supplemental fallback)
KNOWN_BRAND = {
    "primary":    "#0766FF",
    "secondary":  "#00328D",
    "accent":     "#FFA41B",
    "text":       "#00328D",
    "background": "#FFFFFF",
    "font":       "Avantt",
}

# XML namespaces used in PPTX
NS = {
    "a":   "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p":   "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r":   "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "pic": "http://schemas.openxmlformats.org/drawingml/2006/picture",
}

# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def _emu_to_in(emu: int) -> float:
    """Convert EMU integer to inches, rounded to 3 decimal places."""
    return round(emu / 914400, 3)


def _rgb_to_hex(rgb) -> str | None:
    """Convert an RGBColor or (r,g,b) tuple to '#RRGGBB' string."""
    if rgb is None:
        return None
    try:
        if hasattr(rgb, "rgb"):
            rgb = rgb.rgb
        r = int(rgb[0:2], 16) if isinstance(rgb, str) else rgb[0]
        g = int(rgb[2:4], 16) if isinstance(rgb, str) else rgb[1]
        b = int(rgb[4:6], 16) if isinstance(rgb, str) else rgb[2]
        return f"#{r:02X}{g:02X}{b:02X}"
    except Exception:
        return None


def _color_from_element(elem) -> str | None:
    """
    Try to extract a solid RGB colour from a DrawingML element (spPr, solidFill, etc.)
    Returns '#RRGGBB' or None.
    """
    try:
        srgb = elem.find(".//{%s}srgbClr" % NS["a"])
        if srgb is not None:
            val = srgb.get("val", "")
            if len(val) == 6:
                return f"#{val.upper()}"
        # Theme colour with lumMod/lumOff – too complex; skip
    except Exception:
        pass
    return None


def _is_ellipse(shape) -> bool:
    """Return True if the shape has a preset geometry of 'ellipse'."""
    try:
        sp_pr = shape.element.find(".//{%s}prstGeom" % NS["a"])
        if sp_pr is not None and sp_pr.get("prst") == "ellipse":
            return True
        # For picture shapes check spPr > prstGeom
        prstGeom = shape.element.find(".//{%s}prstGeom" % NS["a"])
        if prstGeom is not None and prstGeom.get("prst") == "ellipse":
            return True
    except Exception:
        pass
    return False


def _get_prstGeom(shape) -> str | None:
    """Return the prst attribute of prstGeom if present, else None."""
    try:
        pg = shape.element.find(".//{%s}prstGeom" % NS["a"])
        if pg is not None:
            return pg.get("prst")
    except Exception:
        pass
    return None


def _shape_fill_hex(shape) -> str | None:
    """Extract solid fill colour hex from a shape's spPr."""
    try:
        sp_pr = shape.element.find("{%s}spPr" % NS["p"])
        if sp_pr is None:
            # Try without namespace (auto-resolved)
            for child in shape.element:
                if child.tag.endswith("}spPr") or child.tag == "spPr":
                    sp_pr = child
                    break
        if sp_pr is None:
            return None
        solid = sp_pr.find(".//{%s}solidFill" % NS["a"])
        if solid is None:
            return None
        return _color_from_element(solid)
    except Exception:
        return None


def _is_background_image(shape) -> bool:
    """Heuristic: shape covers >= 80% of slide area (it's a background)."""
    try:
        w = _emu_to_in(shape.width)
        h = _emu_to_in(shape.height)
        return (w * h) >= (SLIDE_AREA * 0.80)
    except Exception:
        return False


def _has_chart(shape) -> bool:
    """Return True if this shape (or group) contains a chart."""
    try:
        if hasattr(shape, "chart"):
            return True
        # Check XML for chart reference
        tag_str = etree.tostring(shape.element, encoding="unicode")
        return "chartSpace" in tag_str or "c:chart" in tag_str or "/chart" in tag_str
    except Exception:
        return False


def _count_groups(slide) -> int:
    """Count top-level group shapes on a slide."""
    count = 0
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                count += 1
        except Exception:
            pass
    return count


def _sample_text(text: str, max_chars: int = 40) -> str:
    """Return first max_chars of text, stripped of newlines."""
    t = text.replace("\n", " ").replace("\r", " ").strip()
    return t[:max_chars] + ("…" if len(t) > max_chars else "")


# ---------------------------------------------------------------------------
# Text zone extraction
# ---------------------------------------------------------------------------

def _extract_text_zones(slide) -> list[dict]:
    """
    Walk all shapes (including inside groups) and extract text zones.
    Returns list of zone descriptors sorted by top position.
    """
    zones = []
    zone_idx = 0

    def _process_shape(shape):
        nonlocal zone_idx
        try:
            shape_type = shape.shape_type
        except Exception:
            shape_type = None

        if shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for child in shape.shapes:
                    _process_shape(child)
            except Exception:
                pass
            return

        if not hasattr(shape, "text_frame"):
            return

        try:
            tf = shape.text_frame
            full_text = tf.text.strip()
            if not full_text:
                return
        except Exception:
            return

        # Gather font info from all runs
        max_pt = 0.0
        font_name = None
        bold = False
        color_hex = None

        try:
            for para in tf.paragraphs:
                for run in para.runs:
                    try:
                        sz = run.font.size
                        if sz:
                            pt = sz / 12700  # EMU to pt
                            if pt > max_pt:
                                max_pt = pt
                        if run.font.name and not font_name:
                            font_name = run.font.name
                        if run.font.bold:
                            bold = True
                        if run.font.color and run.font.color.type:
                            try:
                                hex_val = _rgb_to_hex(run.font.color.rgb)
                                if hex_val and not color_hex:
                                    color_hex = hex_val
                            except Exception:
                                pass
                    except Exception:
                        pass
        except Exception:
            pass

        try:
            left = _emu_to_in(shape.left)
            top = _emu_to_in(shape.top)
            width = _emu_to_in(shape.width)
            height = _emu_to_in(shape.height)
        except Exception:
            left = top = width = height = 0.0

        area = width * height
        n_paras = len([p for p in tf.paragraphs if p.text.strip()])

        # Infer role
        if max_pt >= 60:
            role = "HIGHLIGHT"
        elif max_pt >= 24 and top < 2.0:
            role = "TITLE"
        elif max_pt >= 14 and top < 3.0:
            role = "SUBTITLE"
        elif (area / SLIDE_AREA) > 0.20 and n_paras > 1:
            role = "BODY"
        elif width < 3.0:
            role = "LABEL"
        else:
            role = "CAPTION"

        zones.append({
            "zone_id": f"ZONE_{zone_idx}",
            "role": role,
            "bounds": [left, top, width, height],
            "max_font_pt": round(max_pt, 1),
            "font_name": font_name or "",
            "bold": bold,
            "color": color_hex or "",
            "sample_text": _sample_text(full_text),
        })
        zone_idx += 1

    for shape in slide.shapes:
        _process_shape(shape)

    zones.sort(key=lambda z: (z["bounds"][1], z["bounds"][0]))
    return zones


# ---------------------------------------------------------------------------
# Image zone extraction
# ---------------------------------------------------------------------------

def _extract_image_zones(slide) -> list[dict]:
    """
    Walk all shapes and identify image / picture shapes.
    Returns list of image zone descriptors.
    """
    zones = []
    img_idx = 0

    def _process_shape(shape):
        nonlocal img_idx
        try:
            shape_type = shape.shape_type
        except Exception:
            shape_type = None

        if shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for child in shape.shapes:
                    _process_shape(child)
            except Exception:
                pass
            return

        is_pic = shape_type == MSO_SHAPE_TYPE.PICTURE
        # Some images appear as FREEFORM or other types with pic: XML
        if not is_pic:
            try:
                tag_str = etree.tostring(shape.element, encoding="unicode")
                if "<p:pic" in tag_str or "<pic:pic" in tag_str:
                    is_pic = True
            except Exception:
                pass

        if not is_pic:
            return

        try:
            left = _emu_to_in(shape.left)
            top = _emu_to_in(shape.top)
            width = _emu_to_in(shape.width)
            height = _emu_to_in(shape.height)
        except Exception:
            left = top = width = height = 0.0

        is_bg = _is_background_image(shape)
        is_circle = _is_ellipse(shape)

        if is_bg:
            role = "BACKGROUND"
            crop = "full"
        elif is_circle:
            role = "HERO_CIRCLE"
            crop = "circle"
        else:
            role = "HERO_RECT"
            crop = "rect"

        # Small decorative images
        if not is_bg and (width < 1.5 or height < 1.5):
            role = "DECO"

        zones.append({
            "zone_id": f"IMG_{img_idx}",
            "role": role,
            "bounds": [left, top, width, height],
            "crop": crop,
            "is_background": is_bg,
        })
        img_idx += 1

    for shape in slide.shapes:
        _process_shape(shape)

    return zones


# ---------------------------------------------------------------------------
# Feature detection
# ---------------------------------------------------------------------------

def _detect_features(slide, text_zones: list[dict], image_zones: list[dict]) -> dict:
    """
    Derive boolean/count features for a slide that drive the use_for inference.
    """
    # Text metrics
    text_density = len(text_zones)
    has_title = any(z["role"] == "TITLE" for z in text_zones)
    has_highlight = any(z["role"] == "HIGHLIGHT" for z in text_zones)

    # Image metrics
    img_count = len(image_zones)
    hero_images_circle = sum(1 for z in image_zones if z["role"] == "HERO_CIRCLE")

    # Geometry scan (shapes only, not pictures)
    roundrect_count = 0
    has_chart = False

    def _scan_shape(shape):
        nonlocal roundrect_count, has_chart
        try:
            shape_type = shape.shape_type
        except Exception:
            shape_type = None

        if shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for child in shape.shapes:
                    _scan_shape(child)
            except Exception:
                pass
            return

        prst = _get_prstGeom(shape)
        if prst == "roundRect":
            roundrect_count += 1

        if _has_chart(shape):
            has_chart = True

    for shape in slide.shapes:
        _scan_shape(shape)

    cards_roundrect = roundrect_count
    has_group_clusters = _count_groups(slide) >= 3

    return {
        "has_title":          has_title,
        "has_highlight":      has_highlight,
        "hero_images_circle": hero_images_circle,
        "cards_roundrect":    cards_roundrect,
        "has_chart":          has_chart,
        "has_group_clusters": has_group_clusters,
        "text_density":       text_density,
        "img_count":          img_count,
    }


def _infer_use_for(features: dict) -> list[str]:
    """Map slide features to a list of semantic use-for tags."""
    uses = []
    h = features["has_highlight"]
    hc = features["hero_images_circle"]
    rr = features["cards_roundrect"]
    chart = features["has_chart"]
    td = features["text_density"]
    ic = features["img_count"]

    if h and hc > 0:
        uses.append("hero")
    if rr >= 4:
        uses.append("cards")
    if chart:
        uses += ["dashboard", "chart-placeholder"]
    if hc >= 4:
        uses.append("team")
    if td <= 3 and hc == 0:
        uses += ["title", "section"]
    if hc == 1 and td >= 3:
        uses += ["hero", "image-text"]
    if td >= 6 and not chart:
        uses += ["content", "two-column"]
    if td <= 2 and ic <= 2:
        uses.append("closing")

    if not uses:
        uses.append("content")

    # Deduplicate while preserving order
    seen = set()
    result = []
    for u in uses:
        if u not in seen:
            seen.add(u)
            result.append(u)
    return result


# ---------------------------------------------------------------------------
# Brand token accumulator
# ---------------------------------------------------------------------------

class BrandAccumulator:
    """Accumulates font / colour statistics across all slides."""

    def __init__(self):
        self.font_counter: Counter = Counter()
        self.large_color_counter: Counter = Counter()   # size >= 20pt
        self.small_color_counter: Counter = Counter()   # size < 14pt
        self.fill_color_counter: Counter = Counter()

    def feed_text_zones(self, zones: list[dict]) -> None:
        for z in zones:
            fn = z.get("font_name", "")
            if fn:
                self.font_counter[fn] += 1
            pt = z.get("max_font_pt", 0)
            col = z.get("color", "")
            if col:
                if pt >= 20:
                    self.large_color_counter[col] += 1
                if pt < 14:
                    self.small_color_counter[col] += 1

    def feed_fill_colors(self, slide) -> None:
        def _scan(shape):
            try:
                shape_type = shape.shape_type
            except Exception:
                shape_type = None

            if shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for child in shape.shapes:
                        _scan(child)
                except Exception:
                    pass
                return

            col = _shape_fill_hex(shape)
            if col:
                self.fill_color_counter[col] += 1

        for shape in slide.shapes:
            _scan(shape)

    def build_tokens(self) -> dict:
        brand_font = (
            self.font_counter.most_common(1)[0][0]
            if self.font_counter
            else KNOWN_BRAND["font"]
        )
        # Override with known Avantt if detected (handles partial matches)
        if not any("Avantt" in fn for fn in self.font_counter):
            brand_font = KNOWN_BRAND["font"]

        title_color = (
            self.large_color_counter.most_common(1)[0][0]
            if self.large_color_counter
            else KNOWN_BRAND["secondary"]
        )
        body_color = (
            self.small_color_counter.most_common(1)[0][0]
            if self.small_color_counter
            else KNOWN_BRAND["text"]
        )

        # Top fill colours (exclude white/near-white)
        fill_colors = [
            c for c, _ in self.fill_color_counter.most_common(20)
            if c.upper() not in ("#FFFFFF", "#FEFEFE", "#F0F0F0", "#000000")
        ]

        # Map to palette slots using known values as anchor
        primary   = _pick_closest(fill_colors, KNOWN_BRAND["primary"],   KNOWN_BRAND["primary"])
        secondary = _pick_closest(fill_colors, KNOWN_BRAND["secondary"],  KNOWN_BRAND["secondary"])
        accent    = _pick_closest(fill_colors, KNOWN_BRAND["accent"],     KNOWN_BRAND["accent"])

        return {
            "palette": {
                "primary":    primary,
                "secondary":  secondary,
                "accent":     accent,
                "text":       body_color,
                "background": "#FFFFFF",
            },
            "typography": {
                "title": {
                    "font":    brand_font,
                    "size_pt": 28,
                    "bold":    True,
                    "color":   KNOWN_BRAND["secondary"],
                },
                "subtitle": {
                    "font":    brand_font,
                    "size_pt": 20,
                    "bold":    True,
                    "color":   KNOWN_BRAND["secondary"],
                },
                "body": {
                    "font":    brand_font,
                    "size_pt": 9,
                    "bold":    False,
                    "color":   body_color,
                },
                "kpi": {
                    "font":    brand_font,
                    "size_pt": 28,
                    "bold":    True,
                    "color":   KNOWN_BRAND["accent"],
                },
                "label": {
                    "font":    brand_font,
                    "size_pt": 8,
                    "bold":    False,
                    "color":   KNOWN_BRAND["primary"],
                },
            },
            "slide_width_in":  SLIDE_W_IN,
            "slide_height_in": SLIDE_H_IN,
        }


def _hex_distance(a: str, b: str) -> float:
    """Euclidean distance in RGB space between two '#RRGGBB' strings."""
    def _parse(h):
        h = h.lstrip("#")
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    try:
        r1, g1, b1 = _parse(a)
        r2, g2, b2 = _parse(b)
        return ((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) ** 0.5
    except Exception:
        return 9999.0


def _pick_closest(candidates: list[str], target: str, fallback: str) -> str:
    """Pick the candidate colour closest to target; fallback if none found."""
    if not candidates:
        return fallback
    # Accept exact match first
    for c in candidates:
        if c.upper() == target.upper():
            return c
    # Accept closest within threshold (< 80 in RGB distance)
    best = min(candidates, key=lambda c: _hex_distance(c, target))
    if _hex_distance(best, target) < 80:
        return best
    # Nothing close enough — use known value
    return fallback


# ---------------------------------------------------------------------------
# Main inspection logic
# ---------------------------------------------------------------------------

def inspect_pptx(pptx_path: str, accum: BrandAccumulator) -> list[dict]:
    """
    Inspect a single PPTX file and return its slide catalogue entries.
    Uses open(path, 'rb') + BytesIO to handle Windows paths with spaces.
    """
    fname = os.path.basename(pptx_path)
    print(f"\n  Opening: {fname}")

    try:
        with open(pptx_path, "rb") as fh:
            data = fh.read()
        prs = Presentation(io.BytesIO(data))
    except PermissionError:
        # File locked by OneDrive — use subprocess shell copy to bypass the lock
        import subprocess, tempfile
        print(f"  [RETRY] PermissionError — shell-copying to temp: {fname}")
        try:
            tmp_path = os.path.join(tempfile.gettempdir(), f"ppmaker_inspect_{os.getpid()}.pptx")
            subprocess.run(
                ["cmd", "/c", "copy", "/Y", pptx_path.replace("/", "\\"), tmp_path.replace("/", "\\")],
                check=True, capture_output=True
            )
            with open(tmp_path, "rb") as fh:
                data = fh.read()
            os.unlink(tmp_path)
            prs = Presentation(io.BytesIO(data))
        except Exception as exc2:
            print(f"  [SKIP] Still cannot open after shell copy: {exc2}")
            return []
    except Exception as exc:
        print(f"  [SKIP] Cannot open {fname}: {exc}")
        return []

    entries = []
    n_slides = len(prs.slides)
    print(f"  Slides: {n_slides}")

    for idx, slide in enumerate(prs.slides):
        slide_id = f"SLIDE{idx + 1:02d}"
        try:
            layout_name = slide.slide_layout.name
        except Exception:
            layout_name = ""

        # Extract zones
        text_zones  = _extract_text_zones(slide)
        image_zones = _extract_image_zones(slide)
        features    = _detect_features(slide, text_zones, image_zones)
        use_for     = _infer_use_for(features)

        # Feed brand accumulator
        accum.feed_text_zones(text_zones)
        accum.feed_fill_colors(slide)

        entry = {
            "layout_id":    slide_id,
            "source_file":  fname,
            "slide_index":  idx,
            "layout_name":  layout_name,
            "features":     features,
            "text_zones":   text_zones,
            "image_zones":  image_zones,
            "use_for":      use_for,
        }
        entries.append(entry)

    return entries


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main():
    # Resolve paths relative to ppmaker/ root
    # Also accept an optional --input argument to point at a pre-copied file
    import argparse
    parser = argparse.ArgumentParser(add_help=False)
    parser.add_argument("--input", default=None,
                        help="Path to a specific .pptx (bypasses example dir scan)")
    args, _ = parser.parse_known_args()

    script_dir   = os.path.dirname(os.path.abspath(__file__))   # tools/
    ppmaker_root = os.path.dirname(script_dir)                   # ppmaker/
    example_dir  = os.path.join(ppmaker_root, "template_padrao", "local", "example")
    output_dir   = script_dir  # tools/

    print("=" * 60)
    print("PPMaker Template Inspector")
    print("=" * 60)

    if args.input:
        pptx_files = [args.input]
        print(f"  Using explicit input: {args.input}")
    else:
        print(f"  Looking for .pptx in: {example_dir}")
        pptx_files = glob.glob(os.path.join(example_dir, "*.pptx"))

        # If all files are OneDrive-locked, fall back to a known temp copy
        temp_fallback = os.path.join(os.environ.get("TEMP", os.environ.get("TMP", "/tmp")),
                                     "ambev.pptx")
        if pptx_files and not os.path.isfile(temp_fallback):
            # pre-create the temp copy via subprocess shell (bypasses OneDrive lock)
            src = pptx_files[0]
            try:
                import subprocess
                subprocess.run(["robocopy", os.path.dirname(src),
                                os.path.dirname(temp_fallback),
                                os.path.basename(src),
                                "/COPY:D", "/R:1", "/W:1"],
                               capture_output=True)
                if not os.path.isfile(temp_fallback):
                    subprocess.run(["xcopy", "/Y", src, temp_fallback],
                                   capture_output=True, shell=True)
            except Exception:
                pass

        if not pptx_files:
            print("  [ERROR] No .pptx files found. Exiting.")
            sys.exit(1)

    print(f"  Found {len(pptx_files)} file(s).")

    catalog_entries: list[dict] = []

    print(f"  Found {len(pptx_files)} file(s).")

    catalog_entries: list[dict] = []
    accum = BrandAccumulator()

    for pptx_path in sorted(pptx_files):
        entries = inspect_pptx(pptx_path, accum)
        catalog_entries.extend(entries)

    # Build output artefacts
    catalog = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "slide_dimensions": {
            "width_in":  SLIDE_W_IN,
            "height_in": SLIDE_H_IN,
        },
        "catalog": catalog_entries,
    }

    brand_tokens = accum.build_tokens()

    # Write layout_catalog.json
    catalog_path = os.path.join(output_dir, "layout_catalog.json")
    with open(catalog_path, "w", encoding="utf-8") as fh:
        json.dump(catalog, fh, indent=2, ensure_ascii=False)
    print(f"\n  Written: {catalog_path}")

    # Write brand_tokens.json
    tokens_path = os.path.join(output_dir, "brand_tokens.json")
    with open(tokens_path, "w", encoding="utf-8") as fh:
        json.dump(brand_tokens, fh, indent=2, ensure_ascii=False)
    print(f"  Written: {tokens_path}")

    # ---- Summary report ----
    print("\n" + "=" * 60)
    print("SUMMARY")
    print("=" * 60)
    print(f"  Total slides catalogued : {len(catalog_entries)}")

    use_for_stats: Counter = Counter()
    for e in catalog_entries:
        for u in e["use_for"]:
            use_for_stats[u] += 1

    print("\n  use_for distribution:")
    for tag, cnt in use_for_stats.most_common():
        print(f"    {tag:<20} {cnt}")

    print("\n  Brand tokens (palette):")
    for k, v in brand_tokens["palette"].items():
        print(f"    {k:<12} {v}")

    print("\n  Typography:")
    for role, info in brand_tokens["typography"].items():
        print(f"    {role:<10} {info['font']} {info['size_pt']}pt "
              f"{'bold' if info['bold'] else '    '} {info['color']}")

    print("\nDone.")


if __name__ == "__main__":
    main()
