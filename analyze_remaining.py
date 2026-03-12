"""
Analyze remaining 2 pptx files by opening them as file handles.
"""
import os
import sys
import glob
from collections import Counter, defaultdict
from pptx import Presentation
import traceback

EXAMPLE_DIR = r"C:\Users\99801471\OneDrive - Anheuser-Busch InBev\My Documents\_Meus Documentos\_IA\ppmaker\template_padrao\local\example"

def rgb_to_hex(rgb):
    if rgb is None:
        return None
    return "#{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2])

def get_color_from_fill(fill):
    try:
        if fill is None:
            return None
        ft = fill.type
        if ft is None:
            return None
        if str(ft) == 'SOLID (1)' or ft == 1:
            fc = fill.fore_color
            if fc is not None:
                try:
                    rgb = fc.rgb
                    return rgb_to_hex(rgb)
                except:
                    pass
    except:
        pass
    return None

files = glob.glob(os.path.join(EXAMPLE_DIR, "*.pptx"))
print(f"Files found: {len(files)}")

all_fill_colors = Counter()
all_text_colors = Counter()
all_fonts = Counter()
all_bg_colors = Counter()

for filepath in files:
    fname = os.path.basename(filepath)
    print(f"\nTrying: {fname}")
    try:
        # Open as binary file handle to bypass path encoding
        with open(filepath, 'rb') as f:
            prs = Presentation(f)

        w = prs.slide_width
        h = prs.slide_height
        w_in = w / 914400
        h_in = h / 914400
        print(f"  Slides: {len(prs.slides)}, Size: {w_in:.2f}\" x {h_in:.2f}\"")
        print(f"  Layouts: {[l.name for l in prs.slide_layouts]}")

        for slide_idx, slide in enumerate(prs.slides):
            bg = slide.background.fill
            try:
                bg_color = get_color_from_fill(bg)
                if bg_color:
                    all_bg_colors[bg_color] += 1
            except:
                pass

            for shape in slide.shapes:
                try:
                    if shape.fill:
                        fc = get_color_from_fill(shape.fill)
                        if fc:
                            all_fill_colors[fc] += 1
                except:
                    pass

                try:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                fn = run.font.name
                                fs = run.font.size
                                if fn:
                                    fs_pt = round(fs / 12700) if fs else None
                                    all_fonts[(fn, fs_pt)] += 1
                                try:
                                    if run.font.color and run.font.color.type is not None:
                                        tc = run.font.color.rgb
                                        all_text_colors[rgb_to_hex(tc)] += 1
                                except:
                                    pass
                except:
                    pass

        print(f"  Top fill colors: {list(all_fill_colors.most_common(5))}")
        print(f"  Top text colors: {list(all_text_colors.most_common(5))}")
        print(f"  Top fonts: {list(all_fonts.most_common(5))}")
    except Exception as e:
        print(f"  ERROR: {e}")

print(f"\n\nCOMBINED ACROSS ALL 4 FILES:")
print(f"Fill colors (top 15 non-white):")
shown = 0
for color, count in all_fill_colors.most_common(30):
    try:
        r = int(color[1:3], 16)
        g = int(color[3:5], 16)
        b = int(color[5:7], 16)
        if r > 240 and g > 240 and b > 240:
            continue
    except:
        pass
    print(f"  {color}: {count}")
    shown += 1
    if shown >= 15:
        break

print(f"\nText colors (top 10):")
for color, count in all_text_colors.most_common(10):
    print(f"  {color}: {count}")

print(f"\nFonts (top 15):")
for (fn, fs), count in all_fonts.most_common(15):
    print(f"  {fn} {fs}pt: {count}")

print(f"\nBackground colors:")
for color, count in all_bg_colors.most_common(5):
    print(f"  {color}: {count}")
