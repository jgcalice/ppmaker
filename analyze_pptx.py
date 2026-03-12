"""
Analyze PowerPoint files to extract design patterns.
"""
import os
import sys
from collections import Counter, defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import traceback

FILES = [
    r"C:\Users\99801471\OneDrive - Anheuser-Busch InBev\My Documents\_Meus Documentos\_IA\ppmaker\template_padrao\local\example\08.25 - Monthly Projetos - Gestão de Ativos + Risco.pptx",
    r"C:\Users\99801471\OneDrive - Anheuser-Busch InBev\My Documents\_Meus Documentos\_IA\ppmaker\template_padrao\local\example\2 - Mensal_27.03 1.pptx",
    r"C:\Users\99801471\OneDrive - Anheuser-Busch InBev\My Documents\_Meus Documentos\_IA\ppmaker\template_padrao\local\example\MPR - Commercial Setembro.pptx",
    r"C:\Users\99801471\OneDrive - Anheuser-Busch InBev\My Documents\_Meus Documentos\_IA\ppmaker\template_padrao\local\example\Onboarding Projetos BR Comodatos e Risco.pptx",
]

def rgb_to_hex(rgb):
    if rgb is None:
        return None
    return "#{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2])

def get_color_from_fill(fill):
    """Extract RGB from a fill object."""
    try:
        if fill is None:
            return None
        from pptx.enum.dml import MSO_THEME_COLOR
        ft = fill.type
        if ft is None:
            return None
        from pptx.util import Pt
        # Solid fill
        if str(ft) == 'SOLID (1)' or ft == 1:
            fc = fill.fore_color
            if fc is not None:
                try:
                    rgb = fc.rgb
                    return rgb_to_hex(rgb)
                except:
                    pass
    except Exception:
        pass
    return None

def analyze_file(filepath):
    print(f"\n{'='*70}")
    print(f"FILE: {os.path.basename(filepath)}")
    print('='*70)

    prs = Presentation(filepath)

    # Slide dimensions
    w = prs.slide_width
    h = prs.slide_height
    w_in = w / 914400
    h_in = h / 914400
    print(f"\nSlide dimensions: {w} x {h} EMU  ({w_in:.2f}\" x {h_in:.2f}\")")
    ratio = w_in / h_in
    print(f"Aspect ratio: {ratio:.3f} ({'16:9' if abs(ratio - 1.778) < 0.05 else '4:3' if abs(ratio - 1.333) < 0.05 else 'other'})")

    # Slide master layouts
    master = prs.slide_master
    print(f"\nSlide master layouts ({len(prs.slide_layouts)}):")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  [{i}] {layout.name}")

    # Theme colors
    print(f"\nTheme colors from slide master:")
    try:
        theme_element = master.element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}theme')
        if theme_element is None:
            # Try to access theme directly
            from pptx.oxml.ns import qn
            theme_element = master.element.find('.//' + qn('a:theme'))
        if theme_element is not None:
            clrScheme = theme_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme')
            if clrScheme is not None:
                for child in clrScheme:
                    tag = child.tag.split('}')[-1]
                    for color_child in child:
                        ctag = color_child.tag.split('}')[-1]
                        val = color_child.get('val') or color_child.get('lastClr') or color_child.get('val')
                        if val:
                            print(f"  {tag}: #{val.upper()}")
    except Exception as e:
        print(f"  (could not read theme: {e})")

    # Collect colors and fonts across all slides
    all_bg_colors = []
    all_fill_colors = []
    all_text_colors = []
    all_fonts = []
    shape_types = Counter()
    shape_positions = defaultdict(list)

    for slide_idx, slide in enumerate(prs.slides):
        # Background
        bg = slide.background
        fill = bg.fill
        try:
            bg_color = get_color_from_fill(fill)
            if bg_color:
                all_bg_colors.append(bg_color)
        except:
            pass

        for shape in slide.shapes:
            # Shape type
            try:
                shape_types[str(shape.shape_type)] += 1
            except:
                pass

            # Shape fill
            try:
                if shape.fill:
                    fc = get_color_from_fill(shape.fill)
                    if fc:
                        all_fill_colors.append(fc)
            except:
                pass

            # Text and fonts
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            # Font name
                            fn = run.font.name
                            fs = run.font.size
                            if fn:
                                fs_pt = round(fs / 12700) if fs else None
                                all_fonts.append((fn, fs_pt))
                            # Text color
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    tc = run.font.color.rgb
                                    all_text_colors.append(rgb_to_hex(tc))
                            except:
                                pass
            except:
                pass

            # Track shape positions for recurring elements
            try:
                key = (round(shape.left / 100000), round(shape.top / 100000))
                shape_positions[key].append(slide_idx)
            except:
                pass

    # Report colors
    print(f"\nBackground colors (unique):")
    bg_counter = Counter(all_bg_colors)
    for color, count in bg_counter.most_common(10):
        print(f"  {color}: {count} slides")

    print(f"\nShape fill colors (top 15):")
    fill_counter = Counter(all_fill_colors)
    for color, count in fill_counter.most_common(15):
        print(f"  {color}: {count} shapes")

    print(f"\nText colors (top 10):")
    text_counter = Counter(all_text_colors)
    for color, count in text_counter.most_common(10):
        print(f"  {color}: {count} runs")

    print(f"\nFonts used (top 15):")
    font_counter = Counter(all_fonts)
    for (fname, fsize), count in font_counter.most_common(15):
        print(f"  {fname} {fsize}pt: {count} runs")

    print(f"\nShape types:")
    for stype, count in shape_types.most_common():
        print(f"  {stype}: {count}")

    # Recurring positions (shapes at same position across multiple slides)
    print(f"\nRecurring shape positions (appearing in 3+ slides):")
    for pos, slides_list in sorted(shape_positions.items(), key=lambda x: -len(x[1])):
        if len(slides_list) >= 3:
            print(f"  Position ~({pos[0]*100}, {pos[1]*100}) EMU: appears in {len(slides_list)} slides")

    return {
        'width': w,
        'height': h,
        'bg_colors': bg_counter,
        'fill_colors': fill_counter,
        'text_colors': text_counter,
        'fonts': font_counter,
        'num_slides': len(prs.slides),
    }

def main():
    all_fill_colors = Counter()
    all_text_colors = Counter()
    all_fonts = Counter()
    all_bg_colors = Counter()
    dimensions = []

    for f in FILES:
        if not os.path.exists(f):
            print(f"WARNING: File not found: {f}")
            continue
        try:
            result = analyze_file(f)
            all_fill_colors += result['fill_colors']
            all_text_colors += result['text_colors']
            all_fonts += result['fonts']
            all_bg_colors += result['bg_colors']
            dimensions.append((result['width'], result['height']))
        except Exception as e:
            print(f"ERROR analyzing {f}: {e}")
            traceback.print_exc()

    print(f"\n\n{'='*70}")
    print("COMBINED DESIGN SYSTEM SUMMARY")
    print('='*70)

    if dimensions:
        w, h = dimensions[0]
        w_in = w / 914400
        h_in = h / 914400
        print(f"\nSlide size: {w} x {h} EMU  ({w_in:.2f}\" x {h_in:.2f}\")")

    print(f"\nTOP BACKGROUND COLORS (all files):")
    for color, count in all_bg_colors.most_common(5):
        print(f"  {color}: {count}")

    print(f"\nTOP FILL COLORS (all files, excluding white/near-white):")
    for color, count in all_fill_colors.most_common(20):
        # Skip near-white
        try:
            r = int(color[1:3], 16)
            g = int(color[3:5], 16)
            b = int(color[5:7], 16)
            if r > 240 and g > 240 and b > 240:
                continue
        except:
            pass
        print(f"  {color}: {count}")
        if color == list(all_fill_colors.keys())[19] if len(all_fill_colors) > 19 else True:
            break

    print(f"\nTOP TEXT COLORS (all files):")
    for color, count in all_text_colors.most_common(10):
        print(f"  {color}: {count}")

    print(f"\nTOP FONTS (all files):")
    for (fname, fsize), count in all_fonts.most_common(20):
        print(f"  {fname} {fsize}pt: {count}")

    # Identify dominant design
    print(f"\n{'='*70}")
    print("DOMINANT DESIGN SYSTEM")
    print('='*70)

    # Find primary color (most used non-white/non-black fill)
    primary = None
    secondary = None
    accent = None
    colors_ranked = []
    for color, count in all_fill_colors.most_common(50):
        try:
            r = int(color[1:3], 16)
            g = int(color[3:5], 16)
            b = int(color[5:7], 16)
            # Skip near-white (>240,>240,>240) and near-black (<20,<20,<20) and near-gray
            if r > 230 and g > 230 and b > 230:
                continue
            if r < 20 and g < 20 and b < 20:
                continue
            colors_ranked.append((color, count, r, g, b))
        except:
            pass

    if colors_ranked:
        primary = colors_ranked[0][0]
        print(f"\nPrimary color (most used non-white/black): {primary}")
    if len(colors_ranked) > 1:
        secondary = colors_ranked[1][0]
        print(f"Secondary color: {secondary}")
    if len(colors_ranked) > 2:
        accent = colors_ranked[2][0]
        print(f"Accent color: {accent}")

    # Background
    bg = "#FFFFFF"
    if all_bg_colors:
        bg = all_bg_colors.most_common(1)[0][0]
    print(f"Background color: {bg}")

    # Text color
    text_color = "#000000"
    if all_text_colors:
        text_color = all_text_colors.most_common(1)[0][0]
    print(f"Text color: {text_color}")

    # Title font
    title_font = "Calibri"
    body_font = "Calibri"
    font_name_counter = Counter()
    for (fname, fsize), count in all_fonts.items():
        if fname:
            font_name_counter[fname] += count
    if font_name_counter:
        top_fonts = font_name_counter.most_common(5)
        print(f"\nTop font names: {[f[0] for f in top_fonts]}")
        title_font = top_fonts[0][0]
        body_font = top_fonts[0][0]

    print(f"Title font: {title_font}")
    print(f"Body font: {body_font}")

    print(f"\n--- Use these values for template creation ---")
    print(f"PRIMARY: {primary or '#003366'}")
    print(f"SECONDARY: {secondary or '#F7941D'}")
    print(f"ACCENT: {accent or '#E31837'}")
    print(f"BACKGROUND: {bg}")
    print(f"TEXT: {text_color}")
    print(f"FONT_TITLE: {title_font}")
    print(f"FONT_BODY: {body_font}")

if __name__ == "__main__":
    main()
