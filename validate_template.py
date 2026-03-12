"""Final validation of local-corporate template."""
import json
from pptx import Presentation

# Validate PPTX
pptx_path = r"C:\Users\99801471\OneDrive - Anheuser-Busch InBev\My Documents\_Meus Documentos\_IA\ppmaker\template_padrao\local\local-corporate.pptx"
prs = Presentation(pptx_path)
print("=== PPTX Validation ===")
print(f"Slides: {len(prs.slides)}")
print(f"Slide size: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")
print(f"Slide layouts available: {len(prs.slide_layouts)}")
for i, slide in enumerate(prs.slides):
    print(f"  Slide {i}: {len(slide.shapes)} shapes")

# Validate JSON
json_path = r"C:\Users\99801471\OneDrive - Anheuser-Busch InBev\My Documents\_Meus Documentos\_IA\ppmaker\template_padrao\local\local-corporate.json"
with open(json_path, "r", encoding="utf-8") as f:
    meta = json.load(f)
print()
print("=== JSON Metadata ===")
print(f"ID:         {meta['id']}")
print(f"Name:       {meta['name']}")
print(f"Scope:      {meta['scope']}")
print(f"Palette:    {meta['palette']}")
print(f"Layouts:    {meta['layouts']}")
print(f"Font title: {meta['font_title']}")
print(f"Font body:  {meta['font_body']}")
print()
print("All files created and validated successfully!")
