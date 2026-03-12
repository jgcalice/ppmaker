"""
Create local-corporate.pptx template based on extracted design system.

Extracted Design System:
  PRIMARY:    #0766FF  (blue - most common fill, used for header bars, accents)
  SECONDARY:  #00328D  (dark navy - background color, dominant text color)
  ACCENT:     #FFA41B  (orange/amber - accent color)
  BACKGROUND: #FFFFFF  (white - slide background)
  TEXT:       #00328D  (dark navy - primary text color)
  DARK_TEXT:  #000000  (black - secondary text)
  LIGHT_TEXT: #FFFFFF  (white - text on dark backgrounds)
  FONT:       Avantt   (primary font, fallback to Calibri)

  Slide size: 12192000 x 6858000 EMU (13.33" x 7.50" = 16:9)

  Structural patterns:
  - Shapes at position (0,0): appears on virtually every slide (full-slide bg element)
  - Shapes at position (200,200) EMU: appears on 21 slides (small logo/brand mark)
  - Header bar at top: colored rectangle spanning width
  - Footer at bottom: thin line or text area
  - Title with dark navy (#00328D) text
  - Content area with light gray or white background
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# Brand colors (extracted from files)
PRIMARY     = RGBColor(0x07, 0x66, 0xFF)   # #0766FF  blue
SECONDARY   = RGBColor(0x00, 0x32, 0x8D)   # #00328D  dark navy
ACCENT      = RGBColor(0xFF, 0xA4, 0x1B)   # #FFA41B  amber/orange
BG_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF  white
BG_LIGHT    = RGBColor(0xEF, 0xEF, 0xEF)   # #EFEFEF  light gray
TEXT_DARK   = RGBColor(0x00, 0x32, 0x8D)   # #00328D  navy text
TEXT_BLACK  = RGBColor(0x00, 0x00, 0x00)   # #000000
TEXT_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
TEXT_GRAY   = RGBColor(0x36, 0x37, 0x39)   # #363739

# Slide dimensions (16:9)
SLIDE_W = 12192000   # 13.33 inches
SLIDE_H = 6858000    # 7.50 inches

# Font - use Calibri as fallback since Avantt may not be available system-wide
# We'll specify Avantt as the desired font and python-pptx will use it if available
FONT_NAME = "Calibri"  # Calibri is safe; Avantt specified in runs for environments that have it


def set_shape_fill_solid(shape, rgb_color):
    """Set a shape's fill to a solid color."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color


def set_shape_no_fill(shape):
    """Set a shape to have no fill (transparent)."""
    shape.fill.background()


def set_shape_no_line(shape):
    """Remove the outline from a shape."""
    shape.line.fill.background()


def add_text_to_shape(shape, text, font_size, bold=False, color=None, align=PP_ALIGN.LEFT, font_name=None):
    """Set text in a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name or FONT_NAME
    if color:
        run.font.color.rgb = color


def set_text_frame_margin(shape, left=0, top=0, right=0, bottom=0):
    """Set text frame internal margins."""
    tf = shape.text_frame
    tf.margin_left = left
    tf.margin_top = top
    tf.margin_right = right
    tf.margin_bottom = bottom


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape to a slide."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    if fill_color:
        set_shape_fill_solid(shape, fill_color)
    else:
        set_shape_no_fill(shape)
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        set_shape_no_line(shape)
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT, font_name=None):
    """Add a text box to a slide."""
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    if text:
        add_text_to_shape(txBox, text, font_size, bold, color, align, font_name)
    return txBox


def create_template():
    output_path = r"C:\Users\99801471\OneDrive - Anheuser-Busch InBev\My Documents\_Meus Documentos\_IA\ppmaker\template_padrao\local\local-corporate.pptx"

    prs = Presentation()

    # Set slide dimensions to 16:9 (matching examples)
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    # Helper dimensions
    W = SLIDE_W
    H = SLIDE_H

    # Header bar height: ~8% of slide height
    HEADER_H = int(H * 0.085)   # ~582600 EMU = about 0.64"
    # Footer height
    FOOTER_H = int(H * 0.06)    # ~411480 EMU = about 0.45"
    FOOTER_TOP = H - FOOTER_H

    # Column widths
    COL1_W = int(W * 0.5)
    COL2_W = W - COL1_W - int(W * 0.02)
    COL2_L = COL1_W + int(W * 0.02)

    # Content area (between header and footer)
    CONTENT_TOP = HEADER_H + int(H * 0.02)
    CONTENT_H = FOOTER_TOP - CONTENT_TOP - int(H * 0.02)
    CONTENT_L = int(W * 0.04)
    CONTENT_W = W - 2 * CONTENT_L

    # -------------------------------------------------------
    # SLIDE 1: Title layout (cover slide)
    # -------------------------------------------------------
    slide_layout = prs.slide_layouts[0]  # Use blank layout as base
    slide = prs.slides.add_slide(slide_layout)

    # Full background
    bg_shape = add_rect(slide, 0, 0, W, H, fill_color=SECONDARY)

    # Blue accent bar on left side (~35% width)
    accent_w = int(W * 0.35)
    # Actually use a top accent band instead, matching what we saw (top bar + bottom footer)
    # Large blue rectangle covering top 40%
    top_cover = add_rect(slide, 0, 0, W, int(H * 0.55), fill_color=PRIMARY)

    # Orange accent line at bottom of blue cover
    accent_line = add_rect(slide, 0, int(H * 0.55), W, int(H * 0.008), fill_color=ACCENT)

    # Title text area (in blue section)
    title_box = add_textbox(
        slide,
        int(W * 0.06), int(H * 0.15),
        int(W * 0.85), int(H * 0.30),
        text="Título da Apresentação",
        font_size=36, bold=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT
    )
    set_text_frame_margin(title_box, left=0, top=0)

    # Subtitle text area (in dark section below blue)
    subtitle_box = add_textbox(
        slide,
        int(W * 0.06), int(H * 0.62),
        int(W * 0.60), int(H * 0.20),
        text="Subtítulo ou data",
        font_size=16, bold=False, color=TEXT_WHITE, align=PP_ALIGN.LEFT
    )

    # Footer line
    footer_line = add_rect(slide, 0, int(H * 0.92), W, int(H * 0.005), fill_color=ACCENT)

    # -------------------------------------------------------
    # SLIDE 2: Content layout (standard with title bar)
    # -------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[0])

    # White background
    add_rect(slide2, 0, 0, W, H, fill_color=BG_WHITE)

    # Blue header bar at top
    header = add_rect(slide2, 0, 0, W, HEADER_H, fill_color=PRIMARY)

    # Title text in header
    title2 = add_textbox(
        slide2,
        int(W * 0.03), int(HEADER_H * 0.15),
        int(W * 0.85), int(HEADER_H * 0.75),
        text="Título do Slide",
        font_size=20, bold=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT
    )
    set_text_frame_margin(title2, top=0, bottom=0)

    # Accent left stripe (thin vertical bar)
    add_rect(slide2, 0, HEADER_H, int(W * 0.004), H - HEADER_H - FOOTER_H, fill_color=ACCENT)

    # Content text area
    content2 = add_textbox(
        slide2,
        int(W * 0.04), CONTENT_TOP,
        CONTENT_W, CONTENT_H,
        text="• Ponto de conteúdo 1\n• Ponto de conteúdo 2\n• Ponto de conteúdo 3",
        font_size=13, color=TEXT_DARK, align=PP_ALIGN.LEFT
    )

    # Footer bar
    footer2 = add_rect(slide2, 0, FOOTER_TOP, W, FOOTER_H, fill_color=SECONDARY)

    # Footer text (company/date placeholder)
    footer_text2 = add_textbox(
        slide2,
        int(W * 0.03), FOOTER_TOP + int(FOOTER_H * 0.2),
        int(W * 0.4), int(FOOTER_H * 0.7),
        text="Empresa | Data",
        font_size=8, color=TEXT_WHITE, align=PP_ALIGN.LEFT
    )

    # Page number on right
    page_num2 = add_textbox(
        slide2,
        int(W * 0.92), FOOTER_TOP + int(FOOTER_H * 0.2),
        int(W * 0.06), int(FOOTER_H * 0.7),
        text="1",
        font_size=8, color=TEXT_WHITE, align=PP_ALIGN.RIGHT
    )

    # -------------------------------------------------------
    # SLIDE 3: Two-column layout
    # -------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[0])

    add_rect(slide3, 0, 0, W, H, fill_color=BG_WHITE)

    # Blue header bar
    add_rect(slide3, 0, 0, W, HEADER_H, fill_color=PRIMARY)
    title3 = add_textbox(
        slide3,
        int(W * 0.03), int(HEADER_H * 0.15),
        int(W * 0.85), int(HEADER_H * 0.75),
        text="Título do Slide",
        font_size=20, bold=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT
    )
    set_text_frame_margin(title3, top=0, bottom=0)

    # Divider line between columns
    gap = int(W * 0.01)
    col1_w = int((W - 2 * int(W * 0.04) - gap) / 2)
    col2_l = int(W * 0.04) + col1_w + gap

    # Column 1 label
    add_textbox(
        slide3,
        int(W * 0.04), CONTENT_TOP,
        col1_w, int(H * 0.06),
        text="Coluna 1",
        font_size=13, bold=True, color=TEXT_DARK, align=PP_ALIGN.LEFT
    )
    # Column 1 content
    add_textbox(
        slide3,
        int(W * 0.04), CONTENT_TOP + int(H * 0.08),
        col1_w, CONTENT_H - int(H * 0.08),
        text="• Item 1\n• Item 2\n• Item 3",
        font_size=12, color=TEXT_DARK, align=PP_ALIGN.LEFT
    )

    # Vertical divider line
    add_rect(slide3, col2_l - int(gap/2), CONTENT_TOP, int(W * 0.002), CONTENT_H, fill_color=BG_LIGHT)

    # Column 2 label
    add_textbox(
        slide3,
        col2_l, CONTENT_TOP,
        col1_w, int(H * 0.06),
        text="Coluna 2",
        font_size=13, bold=True, color=TEXT_DARK, align=PP_ALIGN.LEFT
    )
    # Column 2 content
    add_textbox(
        slide3,
        col2_l, CONTENT_TOP + int(H * 0.08),
        col1_w, CONTENT_H - int(H * 0.08),
        text="• Item A\n• Item B\n• Item C",
        font_size=12, color=TEXT_DARK, align=PP_ALIGN.LEFT
    )

    # Footer
    add_rect(slide3, 0, FOOTER_TOP, W, FOOTER_H, fill_color=SECONDARY)
    add_textbox(slide3, int(W * 0.03), FOOTER_TOP + int(FOOTER_H * 0.2), int(W * 0.4), int(FOOTER_H * 0.7),
                text="Empresa | Data", font_size=8, color=TEXT_WHITE)

    # -------------------------------------------------------
    # SLIDE 4: Chart placeholder layout
    # -------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[0])

    add_rect(slide4, 0, 0, W, H, fill_color=BG_WHITE)

    # Blue header bar
    add_rect(slide4, 0, 0, W, HEADER_H, fill_color=PRIMARY)
    title4 = add_textbox(
        slide4,
        int(W * 0.03), int(HEADER_H * 0.15),
        int(W * 0.85), int(HEADER_H * 0.75),
        text="Título do Gráfico",
        font_size=20, bold=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT
    )
    set_text_frame_margin(title4, top=0, bottom=0)

    # Chart placeholder area with light gray background
    chart_top = CONTENT_TOP
    chart_h = CONTENT_H - int(H * 0.03)
    chart_l = int(W * 0.04)
    chart_w = int(W * 0.92)
    chart_placeholder = add_rect(slide4, chart_l, chart_top, chart_w, chart_h, fill_color=BG_LIGHT)
    chart_placeholder.line.color.rgb = RGBColor(0xC0, 0xC0, 0xC0)
    chart_placeholder.line.width = Pt(0.5)

    # Placeholder text in center
    add_textbox(
        slide4,
        chart_l + int(chart_w * 0.3), chart_top + int(chart_h * 0.4),
        int(chart_w * 0.4), int(chart_h * 0.2),
        text="[ Área para Gráfico ]",
        font_size=14, color=RGBColor(0xA0, 0xA0, 0xA0), align=PP_ALIGN.CENTER
    )

    # Footer
    add_rect(slide4, 0, FOOTER_TOP, W, FOOTER_H, fill_color=SECONDARY)
    add_textbox(slide4, int(W * 0.03), FOOTER_TOP + int(FOOTER_H * 0.2), int(W * 0.4), int(FOOTER_H * 0.7),
                text="Empresa | Data", font_size=8, color=TEXT_WHITE)

    # -------------------------------------------------------
    # SLIDE 5: Image + Text layout
    # -------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[0])

    add_rect(slide5, 0, 0, W, H, fill_color=BG_WHITE)

    # Blue header bar
    add_rect(slide5, 0, 0, W, HEADER_H, fill_color=PRIMARY)
    title5 = add_textbox(
        slide5,
        int(W * 0.03), int(HEADER_H * 0.15),
        int(W * 0.85), int(HEADER_H * 0.75),
        text="Título do Slide",
        font_size=20, bold=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT
    )
    set_text_frame_margin(title5, top=0, bottom=0)

    # Image area (left half, light gray placeholder)
    img_l = int(W * 0.03)
    img_t = CONTENT_TOP
    img_w = int(W * 0.44)
    img_h = CONTENT_H
    img_placeholder = add_rect(slide5, img_l, img_t, img_w, img_h, fill_color=BG_LIGHT)
    img_placeholder.line.color.rgb = RGBColor(0xC0, 0xC0, 0xC0)
    img_placeholder.line.width = Pt(0.5)
    add_textbox(
        slide5,
        img_l + int(img_w * 0.25), img_t + int(img_h * 0.45),
        int(img_w * 0.5), int(img_h * 0.15),
        text="[ Imagem ]",
        font_size=13, color=RGBColor(0xA0, 0xA0, 0xA0), align=PP_ALIGN.CENTER
    )

    # Text area (right half)
    text_l = int(W * 0.50)
    text_w = int(W * 0.46)
    add_textbox(
        slide5,
        text_l, CONTENT_TOP,
        text_w, int(H * 0.06),
        text="Subtítulo",
        font_size=16, bold=True, color=TEXT_DARK, align=PP_ALIGN.LEFT
    )
    add_textbox(
        slide5,
        text_l, CONTENT_TOP + int(H * 0.09),
        text_w, CONTENT_H - int(H * 0.09),
        text="• Descrição do conteúdo 1\n• Descrição do conteúdo 2\n• Descrição do conteúdo 3\n\nTexto de apoio adicional pode ser inserido aqui.",
        font_size=12, color=TEXT_DARK, align=PP_ALIGN.LEFT
    )

    # Footer
    add_rect(slide5, 0, FOOTER_TOP, W, FOOTER_H, fill_color=SECONDARY)
    add_textbox(slide5, int(W * 0.03), FOOTER_TOP + int(FOOTER_H * 0.2), int(W * 0.4), int(FOOTER_H * 0.7),
                text="Empresa | Data", font_size=8, color=TEXT_WHITE)

    # -------------------------------------------------------
    # SLIDE 6: Closing layout
    # -------------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[0])

    # Navy full background
    add_rect(slide6, 0, 0, W, H, fill_color=SECONDARY)

    # Blue accent block (top half)
    add_rect(slide6, 0, 0, W, int(H * 0.5), fill_color=PRIMARY)

    # Orange accent line
    add_rect(slide6, 0, int(H * 0.5), W, int(H * 0.008), fill_color=ACCENT)

    # Main closing text
    add_textbox(
        slide6,
        int(W * 0.15), int(H * 0.18),
        int(W * 0.70), int(H * 0.25),
        text="Obrigado!",
        font_size=48, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER
    )

    # Subtitle
    add_textbox(
        slide6,
        int(W * 0.15), int(H * 0.44),
        int(W * 0.70), int(H * 0.12),
        text="",
        font_size=16, color=TEXT_WHITE, align=PP_ALIGN.CENTER
    )

    # Contact info area
    add_textbox(
        slide6,
        int(W * 0.10), int(H * 0.60),
        int(W * 0.80), int(H * 0.20),
        text="Nome | Email | Telefone",
        font_size=14, color=TEXT_WHITE, align=PP_ALIGN.CENTER
    )

    # Bottom accent line
    add_rect(slide6, 0, int(H * 0.92), W, int(H * 0.005), fill_color=ACCENT)

    # Save
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Template saved to: {output_path}")

    # Validate
    print("\nValidating template...")
    prs2 = Presentation(output_path)
    print(f"  Slides: {len(prs2.slides)}")
    print(f"  Slide size: {prs2.slide_width.inches:.2f}\" x {prs2.slide_height.inches:.2f}\"")
    print(f"  Slide layouts available: {len(prs2.slide_layouts)}")

    print("\nSlide summary:")
    slide_names = [
        "0: title",
        "1: content",
        "2: two-column",
        "3: chart-placeholder",
        "4: image-text",
        "5: closing"
    ]
    for i, name in enumerate(slide_names):
        print(f"  Slide {name}: {len(prs2.slides[i].shapes)} shapes")

    return output_path


if __name__ == "__main__":
    create_template()
