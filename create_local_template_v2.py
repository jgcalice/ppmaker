"""
Create local-corporate.pptx template based on extracted design system.

Strategy: Create a presentation with 6 slides (one per layout type).
pptx_service.py calls _remove_existing_slides() to clear slides, then
adds new slides using prs.slide_layouts[idx].

Since python-pptx's default blank presentation provides 11 layouts (0-10),
we have indices 0-5 available for: title, content, two-column, chart-placeholder,
image-text, closing.

The service uses _get_layout() which:
  1. Searches layouts by name
  2. Falls back to index via LAYOUT_INDEX_MAP: title->0, content->1, two-column->2,
     chart-placeholder->3, image-text->4, closing->5

So the template just needs to be a valid .pptx with at least 6 layouts.
The actual visual design is applied when slides are added dynamically.

Extracted Design System (from 2 of 4 example files):
  PRIMARY:    #0766FF  (blue - most common fill)
  SECONDARY:  #00328D  (dark navy - background/text color)
  ACCENT:     #FFA41B  (amber/orange)
  BACKGROUND: #FFFFFF  (white slide background)
  TEXT:       #00328D  (navy - most common text color)
  FONT:       Avantt   (primary; Calibri as system fallback)

  Slide size: 12192000 x 6858000 EMU (13.33" x 7.50" = 16:9)

  Structural patterns found:
  - Shapes at (0,0) on virtually every slide = full-slide background element
  - Shapes at (200,200) on 21/35 slides = recurring brand mark/logo
  - Blue header bar (#0766FF) at top of content slides
  - Dark navy footer (#00328D) at bottom
  - Orange accent line (#FFA41B) as separator
  - Body text in navy (#00328D)
  - Green status indicators (#2DA703, #00B050) for status/KPI
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# Brand colors (extracted from example files)
PRIMARY     = RGBColor(0x07, 0x66, 0xFF)   # #0766FF  blue
SECONDARY   = RGBColor(0x00, 0x32, 0x8D)   # #00328D  dark navy
ACCENT      = RGBColor(0xFF, 0xA4, 0x1B)   # #FFA41B  amber/orange
BG_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF  white
BG_LIGHT    = RGBColor(0xEF, 0xEF, 0xEF)   # #EFEFEF  light gray
TEXT_DARK   = RGBColor(0x00, 0x32, 0x8D)   # #00328D  navy text
TEXT_BLACK  = RGBColor(0x00, 0x00, 0x00)   # #000000
TEXT_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
TEXT_GRAY   = RGBColor(0x36, 0x37, 0x39)   # #363739
GREEN_OK    = RGBColor(0x2D, 0xA7, 0x03)   # #2DA703  status green

# Slide dimensions (16:9)
SLIDE_W = 12192000   # 13.33 inches
SLIDE_H = 6858000    # 7.50 inches

# Font - Calibri is available on all Windows/Office systems
FONT_NAME = "Calibri"


def set_fill_solid(shape, rgb_color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color


def set_no_fill(shape):
    shape.fill.background()


def set_no_line(shape):
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_pt=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height))
    )
    if fill_color:
        set_fill_solid(shape, fill_color)
    else:
        set_no_fill(shape)
    if line_color:
        shape.line.color.rgb = line_color
        if line_pt:
            shape.line.width = Pt(line_pt)
    else:
        set_no_line(shape)
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=12,
                bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height))
    )
    if text:
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = FONT_NAME
        if color:
            run.font.color.rgb = color
    return txBox


def create_slide_title(prs):
    """Slide 0: Title/Cover slide."""
    W, H = SLIDE_W, SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Navy full background
    add_rect(slide, 0, 0, W, H, fill_color=SECONDARY)

    # Blue top band (55% height)
    add_rect(slide, 0, 0, W, H * 0.55, fill_color=PRIMARY)

    # Orange accent separator line
    add_rect(slide, 0, H * 0.55, W, H * 0.008, fill_color=ACCENT)

    # Title text (large, in blue area)
    add_textbox(slide,
        W * 0.06, H * 0.15,
        W * 0.85, H * 0.28,
        text="Título da Apresentação",
        font_size=36, bold=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT)

    # Subtitle text (below title, still in blue area)
    add_textbox(slide,
        W * 0.06, H * 0.44,
        W * 0.65, H * 0.10,
        text="Subtítulo | Departamento",
        font_size=14, bold=False, color=TEXT_WHITE, align=PP_ALIGN.LEFT)

    # Date/info (in dark navy area)
    add_textbox(slide,
        W * 0.06, H * 0.63,
        W * 0.50, H * 0.10,
        text="Mês Ano",
        font_size=14, color=TEXT_WHITE, align=PP_ALIGN.LEFT)

    # Bottom orange accent line
    add_rect(slide, 0, H * 0.93, W, H * 0.005, fill_color=ACCENT)


def create_slide_content(prs):
    """Slide 1: Standard content slide with header bar and bullet area."""
    W, H = SLIDE_W, SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    HEADER_H = H * 0.10
    FOOTER_H = H * 0.07
    FOOTER_TOP = H - FOOTER_H
    CONTENT_TOP = HEADER_H + H * 0.03
    CONTENT_H = FOOTER_TOP - CONTENT_TOP - H * 0.02
    MARGIN_L = W * 0.04

    # White background
    add_rect(slide, 0, 0, W, H, fill_color=BG_WHITE)

    # Blue header bar
    add_rect(slide, 0, 0, W, HEADER_H, fill_color=PRIMARY)

    # Title in header
    add_textbox(slide,
        MARGIN_L, HEADER_H * 0.15,
        W * 0.85, HEADER_H * 0.75,
        text="Título do Slide",
        font_size=20, bold=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT)

    # Left accent stripe
    add_rect(slide, 0, HEADER_H, W * 0.004, H - HEADER_H - FOOTER_H, fill_color=ACCENT)

    # Content text area
    add_textbox(slide,
        MARGIN_L + W * 0.01, CONTENT_TOP,
        W * 0.92, CONTENT_H,
        text="• Ponto principal 1\n\n• Ponto principal 2\n\n• Ponto principal 3",
        font_size=13, color=TEXT_DARK, align=PP_ALIGN.LEFT)

    # Navy footer bar
    add_rect(slide, 0, FOOTER_TOP, W, FOOTER_H, fill_color=SECONDARY)

    # Footer text
    add_textbox(slide,
        MARGIN_L, FOOTER_TOP + FOOTER_H * 0.2,
        W * 0.4, FOOTER_H * 0.65,
        text="Empresa | Data",
        font_size=8, color=TEXT_WHITE, align=PP_ALIGN.LEFT)

    add_textbox(slide,
        W * 0.92, FOOTER_TOP + FOOTER_H * 0.2,
        W * 0.06, FOOTER_H * 0.65,
        text="01",
        font_size=8, color=TEXT_WHITE, align=PP_ALIGN.RIGHT)


def create_slide_two_column(prs):
    """Slide 2: Two-column layout."""
    W, H = SLIDE_W, SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    HEADER_H = H * 0.10
    FOOTER_H = H * 0.07
    FOOTER_TOP = H - FOOTER_H
    CONTENT_TOP = HEADER_H + H * 0.03
    CONTENT_H = FOOTER_TOP - CONTENT_TOP - H * 0.02
    MARGIN_L = W * 0.04

    add_rect(slide, 0, 0, W, H, fill_color=BG_WHITE)

    # Header
    add_rect(slide, 0, 0, W, HEADER_H, fill_color=PRIMARY)
    add_textbox(slide,
        MARGIN_L, HEADER_H * 0.15, W * 0.85, HEADER_H * 0.75,
        text="Título do Slide",
        font_size=20, bold=True, color=TEXT_WHITE)

    # Two columns
    GAP = W * 0.02
    COL_W = (W - 2 * MARGIN_L - GAP) / 2
    COL1_L = MARGIN_L
    COL2_L = MARGIN_L + COL_W + GAP

    # Column 1 header
    add_rect(slide, COL1_L, CONTENT_TOP, COL_W, H * 0.055, fill_color=PRIMARY)
    add_textbox(slide, COL1_L + W * 0.01, CONTENT_TOP + H * 0.008,
        COL_W - W * 0.02, H * 0.04,
        text="Coluna 1", font_size=12, bold=True, color=TEXT_WHITE)

    # Column 1 content
    add_textbox(slide,
        COL1_L, CONTENT_TOP + H * 0.07,
        COL_W, CONTENT_H - H * 0.07,
        text="• Item 1\n• Item 2\n• Item 3\n• Item 4",
        font_size=12, color=TEXT_DARK)

    # Vertical divider
    add_rect(slide, COL2_L - GAP / 2, CONTENT_TOP, W * 0.002, CONTENT_H, fill_color=BG_LIGHT)

    # Column 2 header
    add_rect(slide, COL2_L, CONTENT_TOP, COL_W, H * 0.055, fill_color=SECONDARY)
    add_textbox(slide, COL2_L + W * 0.01, CONTENT_TOP + H * 0.008,
        COL_W - W * 0.02, H * 0.04,
        text="Coluna 2", font_size=12, bold=True, color=TEXT_WHITE)

    # Column 2 content
    add_textbox(slide,
        COL2_L, CONTENT_TOP + H * 0.07,
        COL_W, CONTENT_H - H * 0.07,
        text="• Item A\n• Item B\n• Item C\n• Item D",
        font_size=12, color=TEXT_DARK)

    # Footer
    add_rect(slide, 0, FOOTER_TOP, W, FOOTER_H, fill_color=SECONDARY)
    add_textbox(slide, MARGIN_L, FOOTER_TOP + FOOTER_H * 0.2, W * 0.4, FOOTER_H * 0.65,
        text="Empresa | Data", font_size=8, color=TEXT_WHITE)


def create_slide_chart_placeholder(prs):
    """Slide 3: Chart placeholder slide."""
    W, H = SLIDE_W, SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    HEADER_H = H * 0.10
    FOOTER_H = H * 0.07
    FOOTER_TOP = H - FOOTER_H
    CONTENT_TOP = HEADER_H + H * 0.03
    CONTENT_H = FOOTER_TOP - CONTENT_TOP - H * 0.02
    MARGIN_L = W * 0.04

    add_rect(slide, 0, 0, W, H, fill_color=BG_WHITE)

    # Header
    add_rect(slide, 0, 0, W, HEADER_H, fill_color=PRIMARY)
    add_textbox(slide, MARGIN_L, HEADER_H * 0.15, W * 0.85, HEADER_H * 0.75,
        text="Título do Gráfico",
        font_size=20, bold=True, color=TEXT_WHITE)

    # Optional subtitle/source text area
    add_textbox(slide, MARGIN_L, CONTENT_TOP, W * 0.6, H * 0.05,
        text="Fonte: | Período:",
        font_size=9, color=TEXT_GRAY)

    # Chart placeholder area
    chart_top = CONTENT_TOP + H * 0.06
    chart_h = CONTENT_H - H * 0.07
    chart_w = W - 2 * MARGIN_L

    placeholder = add_rect(slide, MARGIN_L, chart_top, chart_w, chart_h,
        fill_color=BG_LIGHT,
        line_color=RGBColor(0xC0, 0xC0, 0xC0), line_pt=0.5)

    add_textbox(slide,
        MARGIN_L + chart_w * 0.3, chart_top + chart_h * 0.42,
        chart_w * 0.4, chart_h * 0.16,
        text="[ Área para Gráfico ]",
        font_size=14, color=RGBColor(0xA0, 0xA0, 0xA0),
        align=PP_ALIGN.CENTER, italic=True)

    # Footer
    add_rect(slide, 0, FOOTER_TOP, W, FOOTER_H, fill_color=SECONDARY)
    add_textbox(slide, MARGIN_L, FOOTER_TOP + FOOTER_H * 0.2, W * 0.4, FOOTER_H * 0.65,
        text="Empresa | Data", font_size=8, color=TEXT_WHITE)


def create_slide_image_text(prs):
    """Slide 4: Image left + text right layout."""
    W, H = SLIDE_W, SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    HEADER_H = H * 0.10
    FOOTER_H = H * 0.07
    FOOTER_TOP = H - FOOTER_H
    CONTENT_TOP = HEADER_H + H * 0.03
    CONTENT_H = FOOTER_TOP - CONTENT_TOP - H * 0.02
    MARGIN_L = W * 0.04

    add_rect(slide, 0, 0, W, H, fill_color=BG_WHITE)

    # Header
    add_rect(slide, 0, 0, W, HEADER_H, fill_color=PRIMARY)
    add_textbox(slide, MARGIN_L, HEADER_H * 0.15, W * 0.85, HEADER_H * 0.75,
        text="Título do Slide",
        font_size=20, bold=True, color=TEXT_WHITE)

    # Image placeholder (left 45%)
    img_w = W * 0.44
    img_h = CONTENT_H
    img_placeholder = add_rect(slide, MARGIN_L, CONTENT_TOP, img_w, img_h,
        fill_color=BG_LIGHT,
        line_color=RGBColor(0xC0, 0xC0, 0xC0), line_pt=0.5)
    add_textbox(slide,
        MARGIN_L + img_w * 0.25, CONTENT_TOP + img_h * 0.43,
        img_w * 0.5, img_h * 0.14,
        text="[ Imagem ]",
        font_size=13, color=RGBColor(0xA0, 0xA0, 0xA0),
        align=PP_ALIGN.CENTER, italic=True)

    # Text area (right 48%)
    text_l = MARGIN_L + img_w + W * 0.02
    text_w = W - text_l - MARGIN_L

    add_textbox(slide, text_l, CONTENT_TOP, text_w, H * 0.07,
        text="Subtítulo da Seção",
        font_size=16, bold=True, color=TEXT_DARK)

    add_textbox(slide,
        text_l, CONTENT_TOP + H * 0.09,
        text_w, CONTENT_H - H * 0.09,
        text="• Descrição 1\n• Descrição 2\n• Descrição 3\n\nTexto adicional de apoio.",
        font_size=12, color=TEXT_DARK)

    # Footer
    add_rect(slide, 0, FOOTER_TOP, W, FOOTER_H, fill_color=SECONDARY)
    add_textbox(slide, MARGIN_L, FOOTER_TOP + FOOTER_H * 0.2, W * 0.4, FOOTER_H * 0.65,
        text="Empresa | Data", font_size=8, color=TEXT_WHITE)


def create_slide_closing(prs):
    """Slide 5: Closing/Thank you slide."""
    W, H = SLIDE_W, SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Navy full background
    add_rect(slide, 0, 0, W, H, fill_color=SECONDARY)

    # Blue top block
    add_rect(slide, 0, 0, W, H * 0.52, fill_color=PRIMARY)

    # Orange accent separator
    add_rect(slide, 0, H * 0.52, W, H * 0.008, fill_color=ACCENT)

    # Main text: "Obrigado!"
    add_textbox(slide,
        W * 0.15, H * 0.15,
        W * 0.70, H * 0.27,
        text="Obrigado!",
        font_size=48, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    # Subtitle placeholder
    add_textbox(slide,
        W * 0.20, H * 0.43,
        W * 0.60, H * 0.09,
        text="",
        font_size=16, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    # Contact/info area
    add_textbox(slide,
        W * 0.10, H * 0.62,
        W * 0.80, H * 0.16,
        text="Nome | Email | www.empresa.com.br",
        font_size=14, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    # Bottom orange line
    add_rect(slide, 0, H * 0.93, W, H * 0.005, fill_color=ACCENT)


def main():
    output_path = r"C:\Users\99801471\OneDrive - Anheuser-Busch InBev\My Documents\_Meus Documentos\_IA\ppmaker\template_padrao\local\local-corporate.pptx"

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    print("Creating slides...")
    create_slide_title(prs)
    print("  [0] title - done")
    create_slide_content(prs)
    print("  [1] content - done")
    create_slide_two_column(prs)
    print("  [2] two-column - done")
    create_slide_chart_placeholder(prs)
    print("  [3] chart-placeholder - done")
    create_slide_image_text(prs)
    print("  [4] image-text - done")
    create_slide_closing(prs)
    print("  [5] closing - done")

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"\nTemplate saved to: {output_path}")

    # Validate
    print("\nValidating...")
    prs2 = Presentation(output_path)
    print(f"  Slides: {len(prs2.slides)}")
    print(f"  Slide size: {prs2.slide_width.inches:.2f}\" x {prs2.slide_height.inches:.2f}\"")
    print(f"  Layouts available: {len(prs2.slide_layouts)}")
    for i, slide in enumerate(prs2.slides):
        print(f"  Slide {i}: {len(slide.shapes)} shapes")

    print("\nDone! Template is ready.")


if __name__ == "__main__":
    main()
